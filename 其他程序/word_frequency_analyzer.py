#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
词频分析工具 - 增强版 v7.8
修复：配置保存问题、透明背景与轮廓冲突问题、PDF导出问题
修复：PDF格式词云图显示空白问题
修复：PPT报告中缺少词性分析页面的问题
增强：高频词汇页添加词性信息
"""

import os
import re
import json
import pandas as pd
from collections import Counter
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import threading
from PIL import Image, ImageTk
import matplotlib
# 设置matplotlib使用非交互式后端
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from wordcloud import WordCloud, ImageColorGenerator
import numpy as np
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import jieba
import jieba.analyse
import jieba.posseg as pseg  # 用于词性标注
import queue
import time
import sys

# 尝试导入chardet，如果不存在则使用备用方案
try:
    import chardet
    HAS_CHARDET = True
except ImportError:
    HAS_CHARDET = False
    print("注意: chardet模块未安装，将使用内置编码检测")

# 尝试导入matplotlib PDF后端
try:
    import matplotlib.backends.backend_pdf
    HAS_PDF_SUPPORT = True
except ImportError:
    HAS_PDF_SUPPORT = False
    print("注意: matplotlib PDF后端未安装，PDF导出功能不可用")

# 初始化结巴分词
jieba.initialize()

# 词性对照表 - 完整版
POS_MAP = {
    # 基本词性
    'n': '名词', 'v': '动词', 'a': '形容词', 'd': '副词', 'm': '数词',
    'q': '量词', 'r': '代词', 'p': '介词', 'c': '连词', 'u': '助词',
    'xc': '其他虚词', 'w': '标点符号', 
    
    # 专有名词
    'nr': '人名', 'ns': '地名', 'nt': '机构团体', 'nz': '其他专名',
    
    # 时间、方位、处所
    't': '时间词', 'f': '方位词', 's': '处所词',
    
    # 其他词性
    'vn': '动名词', 'an': '形名词', 'eng': '英文', 'x': '非语素字',
    
    # 词性标注中出现的特殊词性
    'l': '习用语', 'i': '成语', 'b': '区别词', 'vq': '趋向动词', 'ad': '副形词',
    
    # 实体命名识别
    'PER': '人名', 'LOC': '地名', 'ORG': '机构名', 'TIME': '时间',
    
    # jieba其他词性标注
    'g': '语素字', 'h': '前接成分', 'k': '后接成分',
    'o': '拟声词', 'e': '叹词', 'y': '语气词', 'z': '状态词',
    'un': '未知词', 'df': '副词性语素', 'dg': '形容词性语素',
    'uv': '助动词', 'ug': '动词性语素', 'uj': '助词性语素',
    'ul': '连词性语素', 'nrfg': '古代人名', 'nrt': '音译人名',
    'vd': '副动词', 'ag': '形容词性语素', 'vg': '动词性语素',
    
    # 更多可能的词性
    'mg': '数词性语素', 'mq': '数量词', 'tg': '时间词性语素',
    'bg': '区别词性语素', 'rr': '人称代词', 'rz': '指示代词',
    'ry': '疑问代词', 'rg': '代词性语素',
}

class WordFrequencyAnalyzer:
    def __init__(self, root):
        self.root = root
        self.root.title("词频分析工具 v7.8 - 带词性分析和词性管理")
        self.root.geometry("1200x800")  # 更大的窗口
        
        # 设置主题色
        self.bg_color = "#f0f0f0"
        self.btn_color = "#4a7a8c"
        self.text_color = "#333333"
        
        self.root.configure(bg=self.bg_color)
        
        # 初始化变量
        self.file_paths = []
        self.stop_words = set()
        self.custom_words = set()
        self.word_freq = {}
        self.word_pos_freq = {}  # 词性频率统计
        self.word_pos_info = {}  # 每个词的词性信息
        self.bg_image = None
        self.bg_mask = None
        self.bg_image_path = None  # 新增：保存背景图片路径
        
        # 词云设置默认值
        self.wordcloud_settings = {
            'generate_wordcloud': True,
            'generate_ppt': True,  # 默认启用PPT生成
            'background_color': 'white',
            'colormap': 'viridis',
            'max_words': 200,
            'max_font_size': 150,
            'min_font_size': 10,
            'width': 800,
            'height': 600,
            'random_state': 42,
            'repeat': False,
            'include_numbers': False,
            'prefer_horizontal': 0.9,
            'scale': 1,
            'background_image': None,
            'image_contour': False,
            'contour_width': 1,
            'contour_color': 'white',
            'output_format': 'png',
            'font_path': None
        }
        
        # 跟踪设置是否已保存
        self.settings_saved = True
        
        # 线程安全队列
        self.task_queue = queue.Queue()
        self.result_queue = queue.Queue()
        
        # 加载默认配置
        self.load_default_config()
        
        # 创建界面
        self.create_widgets()
        
        # 启动任务处理线程
        self.start_task_handler()
        
        # 定期检查结果队列
        self.check_results()
        
        # 绑定关闭事件
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)
    
    def on_closing(self):
        """关闭窗口时的处理"""
        # 保存所有配置
        self.save_all_config()
        self.root.destroy()
    
    def save_all_config(self):
        """保存所有配置"""
        try:
            # 保存词云设置
            self.save_wordcloud_settings()
            
            # 保存用户配置（停用词和自定义词）
            self.save_user_config()
            
            # 保存界面设置
            self.save_ui_config()
            
            print("所有配置已保存")
        except Exception as e:
            print(f"保存配置时出错: {e}")
    
    def save_ui_config(self):
        """保存界面设置（保存路径等）"""
        try:
            ui_config_file = os.path.join(self.config_dir, "ui_config.json")
            ui_config = {
                "result_dir": self.result_dir_var.get() if hasattr(self, 'result_dir_var') else "",
                "cloud_dir": self.cloud_dir_var.get() if hasattr(self, 'cloud_dir_var') else "",
                "last_tab": self.notebook.index(self.notebook.select()) if hasattr(self, 'notebook') else 0
            }
            
            with open(ui_config_file, 'w', encoding='utf-8') as f:
                json.dump(ui_config, f, ensure_ascii=False, indent=2)
            
            print("界面配置已保存")
        except Exception as e:
            print(f"保存界面配置时出错: {e}")
    
    def load_ui_config(self):
        """加载界面设置"""
        try:
            ui_config_file = os.path.join(self.config_dir, "ui_config.json")
            if os.path.exists(ui_config_file):
                with open(ui_config_file, 'r', encoding='utf-8') as f:
                    ui_config = json.load(f)
                
                # 设置保存路径
                if hasattr(self, 'result_dir_var') and "result_dir" in ui_config:
                    self.result_dir_var.set(ui_config["result_dir"])
                
                if hasattr(self, 'cloud_dir_var') and "cloud_dir" in ui_config:
                    self.cloud_dir_var.set(ui_config["cloud_dir"])
                
                # 设置上次选择的标签页
                if hasattr(self, 'notebook') and "last_tab" in ui_config:
                    try:
                        self.notebook.select(ui_config["last_tab"])
                    except:
                        pass
                
                print("界面配置已加载")
        except Exception as e:
            print(f"加载界面配置时出错: {e}")
    
    def start_task_handler(self):
        """启动任务处理线程"""
        self.handler_thread = threading.Thread(target=self.process_tasks)
        self.handler_thread.daemon = True
        self.handler_thread.start()
    
    def process_tasks(self):
        """处理任务队列中的任务"""
        while True:
            try:
                task = self.task_queue.get(timeout=0.1)
                if task['type'] == 'analysis':
                    self.perform_analysis_task(task)
                elif task['type'] == 'wordcloud':
                    self.generate_wordcloud_task(task)
                self.task_queue.task_done()
            except queue.Empty:
                continue
            except Exception as e:
                print(f"任务处理错误: {e}")
    
    def check_results(self):
        """定期检查结果队列"""
        try:
            while not self.result_queue.empty():
                result = self.result_queue.get_nowait()
                self.handle_result(result)
        except queue.Empty:
            pass
        
        # 100毫秒后再次检查
        self.root.after(100, self.check_results)
    
    def handle_result(self, result):
        """处理结果"""
        if result['type'] == 'analysis_complete':
            self.display_results(result['data'])
            self.status_var.set("分析完成！")
            
            # 保存结果
            self.save_results(result['data']['top_n'])
            
            # 将词云生成任务加入队列
            if self.wordcloud_settings['generate_wordcloud']:
                wordcloud_task = {
                    'type': 'wordcloud',
                    'word_freq': result['data']['word_freq'],
                    'top_n': result['data']['top_n'],
                    'settings': self.wordcloud_settings.copy()
                }
                self.task_queue.put(wordcloud_task)
            
        elif result['type'] == 'wordcloud_complete':
            image_path = result['data'].get('image_path')
            ppt_path = result['data'].get('ppt_path')
            png_preview_path = result['data'].get('png_preview_path')  # 新增：PNG预览文件路径
            
            # 显示词云图（优先使用PNG预览文件）
            display_path = png_preview_path if png_preview_path else image_path
            if display_path and os.path.exists(display_path):
                self.show_wordcloud_image(display_path)
            
            if image_path:
                # 在结果中显示主文件路径
                if os.path.splitext(image_path)[1].lower() == '.pdf':
                    self.result_text.insert(tk.END, f"\n📄 PDF词云图已保存至:\n{image_path}\n")
                else:
                    self.result_text.insert(tk.END, f"\n☁️ 词云图已保存至:\n{image_path}\n")
            
            if ppt_path:
                self.result_text.insert(tk.END, f"\n📊 PPT报告已保存至:\n{ppt_path}\n")
                # 更新界面显示PPT路径
                self.update_ppt_path_display(ppt_path)
            
        elif result['type'] == 'error':
            self.status_var.set("分析出错")
            messagebox.showerror("错误", result['data'])
            self.analyze_btn.config(state='normal', text="开始分析")
    
    def update_ppt_path_display(self, ppt_path):
        """更新界面显示PPT路径"""
        # 在结果文本中显示PPT路径
        if hasattr(self, 'result_text'):
            self.result_text.insert(tk.END, f"\n📋 PPT报告路径已在界面显示\n")
        
        # 也可以在状态栏显示
        self.status_var.set(f"PPT报告已生成: {os.path.basename(ppt_path)}")
    
    def load_default_config(self):
        """加载默认停用词和自定义词"""
        # 默认停用词
        default_stopwords = ["的", "了", "和", "是", "在", "我", "有", "他", "这", "那", 
                            "就", "人", "都", "而", "及", "与", "等", "或", "日", "月",
                            "年", "中", "对", "于", "为", "上", "也", "你", "我", "它",
                            "a", "an", "the", "and", "or", "but", "in", "on", "at", 
                            "to", "for", "of", "with", "by", "from", "as", "is", "are",
                            "was", "were", "be", "been", "being", "have", "has", "had",
                            "do", "does", "did", "will", "would", "shall", "should",
                            "can", "could", "may", "might", "must"]
        
        self.stop_words = set(default_stopwords)
        
        # 保存配置路径
        self.config_dir = os.path.expanduser("~/.word_freq_analyzer")
        os.makedirs(self.config_dir, exist_ok=True)
        
        # 词性映射文件路径
        self.pos_map_file = os.path.join(self.config_dir, "pos_map.json")
        
        # 词云设置文件路径
        self.wordcloud_config_file = os.path.join(self.config_dir, "wordcloud_config.json")
        
        # 尝试加载用户配置
        self.load_user_config()
    
    def load_user_config(self):
        """加载用户配置"""
        config_file = os.path.join(self.config_dir, "config.json")
        if os.path.exists(config_file):
            try:
                with open(config_file, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    if "stop_words" in config:
                        self.stop_words.update(config["stop_words"])
                    if "custom_words" in config:
                        self.custom_words.update(config["custom_words"])
            except Exception as e:
                print(f"加载用户配置时出错: {e}")
        
        # 加载用户自定义词性映射
        if os.path.exists(self.pos_map_file):
            try:
                with open(self.pos_map_file, 'r', encoding='utf-8') as f:
                    user_pos_map = json.load(f)
                    # 更新POS_MAP，但保留原有的
                    for key, value in user_pos_map.items():
                        if key not in POS_MAP:  # 只添加用户新增的词性
                            POS_MAP[key] = value
            except Exception as e:
                print(f"加载用户词性映射时出错: {e}")
        
        # 加载词云设置
        if os.path.exists(self.wordcloud_config_file):
            try:
                with open(self.wordcloud_config_file, 'r', encoding='utf-8') as f:
                    wordcloud_config = json.load(f)
                    self.wordcloud_settings.update(wordcloud_config)
                    
                    # 如果保存了背景图片路径，加载图片
                    if 'background_image' in wordcloud_config and wordcloud_config['background_image']:
                        try:
                            self.bg_image_path = wordcloud_config['background_image']
                            if os.path.exists(self.bg_image_path):
                                self.bg_image = Image.open(self.bg_image_path)
                                self.bg_mask = np.array(self.bg_image)
                        except Exception as e:
                            print(f"加载背景图片时出错: {e}")
                            self.bg_image = None
                            self.bg_mask = None
                            self.bg_image_path = None
            except Exception as e:
                print(f"加载词云配置时出错: {e}")
        
        # 加载界面配置（必须在界面创建后调用）
        # 这个会在create_widgets之后调用
    
    def save_user_config(self):
        """保存用户配置"""
        config_file = os.path.join(self.config_dir, "config.json")
        config = {
            "stop_words": list(self.stop_words),
            "custom_words": list(self.custom_words)
        }
        try:
            with open(config_file, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"保存用户配置时出错: {e}")
    
    def save_wordcloud_settings(self):
        """保存词云配置"""
        try:
            # 首先更新词云设置字典
            self.update_wordcloud_settings_from_ui()
            
            # 更新背景图片路径
            self.wordcloud_settings['background_image'] = self.bg_image_path
            
            with open(self.wordcloud_config_file, 'w', encoding='utf-8') as f:
                json.dump(self.wordcloud_settings, f, ensure_ascii=False, indent=2)
            
            self.settings_saved = True
            return True
        except Exception as e:
            print(f"保存词云配置时出错: {e}")
            return False
    
    def update_wordcloud_settings_from_ui(self):
        """从UI控件更新词云设置字典"""
        if hasattr(self, 'generate_wordcloud_var'):
            self.wordcloud_settings['generate_wordcloud'] = self.generate_wordcloud_var.get()
        if hasattr(self, 'generate_ppt_var'):
            self.wordcloud_settings['generate_ppt'] = self.generate_ppt_var.get()
        if hasattr(self, 'bg_color_var'):
            self.wordcloud_settings['background_color'] = self.bg_color_var.get()
        if hasattr(self, 'colormap_var'):
            self.wordcloud_settings['colormap'] = self.colormap_var.get()
        if hasattr(self, 'max_words_var'):
            self.wordcloud_settings['max_words'] = self.max_words_var.get()
        if hasattr(self, 'max_font_size_var'):
            self.wordcloud_settings['max_font_size'] = self.max_font_size_var.get()
        if hasattr(self, 'min_font_size_var'):
            self.wordcloud_settings['min_font_size'] = self.min_font_size_var.get()
        if hasattr(self, 'random_state_var'):
            self.wordcloud_settings['random_state'] = self.random_state_var.get()
        if hasattr(self, 'repeat_var'):
            self.wordcloud_settings['repeat'] = self.repeat_var.get()
        if hasattr(self, 'include_numbers_var'):
            self.wordcloud_settings['include_numbers'] = self.include_numbers_var.get()
        if hasattr(self, 'prefer_horizontal_var'):
            self.wordcloud_settings['prefer_horizontal'] = self.prefer_horizontal_var.get()
        if hasattr(self, 'output_format_var'):
            self.wordcloud_settings['output_format'] = self.output_format_var.get()
        if hasattr(self, 'font_path_var'):
            self.wordcloud_settings['font_path'] = self.font_path_var.get() if self.font_path_var.get() else None
        if hasattr(self, 'image_contour_var'):
            self.wordcloud_settings['image_contour'] = self.image_contour_var.get()
        if hasattr(self, 'contour_width_var'):
            self.wordcloud_settings['contour_width'] = self.contour_width_var.get()
        if hasattr(self, 'contour_color_var'):
            self.wordcloud_settings['contour_color'] = self.contour_color_var.get()
    
    def detect_encoding_simple(self, file_path):
        """简易编码检测（不使用chardet）"""
        # 常见中文编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5', 'utf-16']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    f.read(1024)  # 读取前1KB测试
                return encoding
            except UnicodeDecodeError:
                continue
            except Exception:
                continue
        
        # 如果都不行，尝试latin-1（总是成功但可能乱码）
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                f.read(1024)
            return 'latin-1'
        except:
            return 'utf-8'  # 默认使用utf-8
    
    def detect_encoding(self, file_path):
        """检测文件编码"""
        if HAS_CHARDET:
            try:
                with open(file_path, 'rb') as f:
                    raw_data = f.read(10000)
                    result = chardet.detect(raw_data)
                    encoding = result['encoding']
                    confidence = result['confidence']
                    
                    if confidence < 0.7:
                        return self.detect_encoding_simple(file_path)
                    
                    return encoding or 'utf-8'
            except Exception as e:
                print(f"chardet检测编码时出错: {e}")
                return self.detect_encoding_simple(file_path)
        else:
            return self.detect_encoding_simple(file_path)
    
    def create_widgets(self):
        """创建界面组件 - 使用Notebook标签页布局"""
        # 创建标题
        title_frame = tk.Frame(self.root, bg=self.bg_color, height=60)
        title_frame.pack(fill=tk.X, pady=5)
        title_frame.pack_propagate(False)
        
        title_label = tk.Label(title_frame, text="📊 词频分析工具 v7.8 - 带词性分析和词性管理", 
                              font=("微软雅黑", 24, "bold"),
                              bg=self.bg_color, fg=self.btn_color)
        title_label.pack(expand=True)
        
        # 创建主Notebook（标签页）
        self.notebook = ttk.Notebook(self.root)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=(0, 10))
        
        # 创建6个标签页
        self.create_file_tab()
        self.create_settings_tab()
        self.create_vocab_tab()
        self.create_pos_tab()  # 词性管理标签页
        self.create_wordcloud_tab()  # 新增词云设置标签页
        self.create_result_tab()
        
        # 底部状态栏和操作按钮
        self.create_bottom_panel()
        
        # 加载界面配置（必须在界面创建后调用）
        self.load_ui_config()
    
    def create_file_tab(self):
        """创建文件管理标签页"""
        file_tab = tk.Frame(self.notebook, bg=self.bg_color)
        self.notebook.add(file_tab, text="📁 文件管理")
        
        # 文件选择区域
        file_frame = tk.LabelFrame(file_tab, text="选择分析文档", 
                                  font=("微软雅黑", 12, "bold"),
                                  bg=self.bg_color, fg=self.text_color)
        file_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # 文件列表区域
        list_frame = tk.Frame(file_frame, bg=self.bg_color)
        list_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # 文件列表框
        listbox_frame = tk.Frame(list_frame, bg=self.bg_color)
        listbox_frame.pack(fill=tk.BOTH, expand=True)
        
        self.file_listbox = tk.Listbox(listbox_frame, font=("宋体", 11), selectmode=tk.EXTENDED)
        self.file_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        # 滚动条
        scrollbar = tk.Scrollbar(listbox_frame, command=self.file_listbox.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.file_listbox.config(yscrollcommand=scrollbar.set)
        
        # 按钮区域
        btn_frame = tk.Frame(file_frame, bg=self.bg_color)
        btn_frame.pack(fill=tk.X, padx=10, pady=10)
        
        # 按钮样式
        btn_style = {'font': ("微软雅黑", 10), 'height': 2, 'width': 15}
        
        add_btn = tk.Button(btn_frame, text="添加文档", command=self.add_files,
                           bg=self.btn_color, fg="white", **btn_style)
        add_btn.pack(side=tk.LEFT, padx=5)
        
        remove_btn = tk.Button(btn_frame, text="移除选中", command=self.remove_file,
                              bg="#e74c3c", fg="white", **btn_style)
        remove_btn.pack(side=tk.LEFT, padx=5)
        
        clear_btn = tk.Button(btn_frame, text="清空列表", command=self.clear_file_list,
                             bg="#e67e22", fg="white", **btn_style)
        clear_btn.pack(side=tk.LEFT, padx=5)
        
        # 文件统计信息
        info_frame = tk.Frame(file_frame, bg=self.bg_color)
        info_frame.pack(fill=tk.X, padx=10, pady=(5, 10))
        
        self.file_count_label = tk.Label(info_frame, text="已选择 0 个文档", 
                                        font=("微软雅黑", 10), bg=self.bg_color)
        self.file_count_label.pack()
    
    def clear_file_list(self):
        """清空文件列表"""
        if self.file_paths:
            if messagebox.askyesno("确认", "确定要清空所有文件吗？"):
                self.file_listbox.delete(0, tk.END)
                self.file_paths.clear()
                self.file_count_label.config(text="已选择 0 个文档")
                self.status_var.set("文件列表已清空")
    
    def create_settings_tab(self):
        """创建分析设置标签页"""
        settings_tab = tk.Frame(self.notebook, bg=self.bg_color)
        self.notebook.add(settings_tab, text="⚙️ 分析设置")
        
        # 使用Canvas和Frame实现滚动
        canvas = tk.Canvas(settings_tab, bg=self.bg_color, highlightthickness=0)
        scrollbar = ttk.Scrollbar(settings_tab, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg=self.bg_color)
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        # 布局
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # 分析设置区域
        analysis_frame = tk.LabelFrame(scrollable_frame, text="基本设置", 
                                      font=("微软雅黑", 12, "bold"),
                                      bg=self.bg_color, fg=self.text_color)
        analysis_frame.pack(fill=tk.X, padx=20, pady=20)
        
        # 显示词数设置 - 改进：更人性化的输入框
        tk.Label(analysis_frame, text="显示词数:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        top_n_frame = tk.Frame(analysis_frame, bg=self.bg_color)
        top_n_frame.pack(fill=tk.X, padx=10, pady=5)
        
        # 滑动条
        self.top_n_var = tk.StringVar(value="50")
        top_n_scale = tk.Scale(top_n_frame, from_=1, to=200, orient=tk.HORIZONTAL,
                              variable=self.top_n_var, bg=self.bg_color,
                              length=300)
        top_n_scale.pack(side=tk.LEFT)
        
        # 输入框 - 改进：使用延迟验证
        top_n_entry_frame = tk.Frame(top_n_frame, bg=self.bg_color)
        top_n_entry_frame.pack(side=tk.LEFT, padx=(20, 0))
        
        tk.Label(top_n_entry_frame, text="输入:", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT)
        
        self.top_n_entry = tk.Entry(top_n_entry_frame, 
                                   font=("宋体", 10), width=10, justify='center')
        self.top_n_entry.insert(0, "50")
        self.top_n_entry.pack(side=tk.LEFT, padx=5)
        
        # 绑定事件，使用延时验证
        def validate_top_n(event=None):
            try:
                value = int(self.top_n_entry.get())
                if 1 <= value <= 200:
                    self.top_n_var.set(str(value))
                    top_n_scale.set(value)
                else:
                    messagebox.showwarning("警告", "请输入1-200之间的整数")
                    self.top_n_entry.delete(0, tk.END)
                    self.top_n_entry.insert(0, "50")
                    self.top_n_var.set("50")
                    top_n_scale.set(50)
            except ValueError:
                # 允许空值或正在输入
                pass
        
        # 失去焦点时验证
        self.top_n_entry.bind('<FocusOut>', validate_top_n)
        # 按回车时验证
        self.top_n_entry.bind('<Return>', validate_top_n)
        
        # 编码设置
        tk.Label(analysis_frame, text="文本编码:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(20, 5))
        
        self.encoding_var = tk.StringVar(value="auto")
        encoding_frame = tk.Frame(analysis_frame, bg=self.bg_color)
        encoding_frame.pack(fill=tk.X, padx=10, pady=5)
        
        encodings = ["auto", "utf-8", "gbk", "gb2312", "gb18030", "big5"]
        encoding_combo = ttk.Combobox(encoding_frame, textvariable=self.encoding_var,
                                     values=encodings, state="readonly", 
                                     font=("宋体", 10), width=20)
        encoding_combo.pack(side=tk.LEFT)
        
        # CSV编码设置
        tk.Label(analysis_frame, text="CSV文件编码:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(20, 5))
        
        self.csv_encoding_var = tk.StringVar(value="utf-8-sig")
        csv_encoding_frame = tk.Frame(analysis_frame, bg=self.bg_color)
        csv_encoding_frame.pack(fill=tk.X, padx=10, pady=5)
        
        csv_encodings = ["utf-8-sig", "gbk", "utf-8", "gb2312", "gb18030"]
        csv_encoding_combo = ttk.Combobox(csv_encoding_frame, textvariable=self.csv_encoding_var,
                                         values=csv_encodings, state="readonly", 
                                         font=("宋体", 10), width=20)
        csv_encoding_combo.pack(side=tk.LEFT)
        
        # 词性分析设置
        tk.Label(analysis_frame, text="词性分析:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(20, 5))
        
        self.enable_pos_var = tk.BooleanVar(value=True)
        tk.Checkbutton(analysis_frame, text="启用词性标注", 
                      variable=self.enable_pos_var,
                      bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=5)
        
        # 保存设置区域
        save_frame = tk.LabelFrame(scrollable_frame, text="保存设置", 
                                  font=("微软雅黑", 12, "bold"),
                                  bg=self.bg_color, fg=self.text_color)
        save_frame.pack(fill=tk.X, padx=20, pady=(0, 20))
        
        # 词频结果保存
        tk.Label(save_frame, text="词频结果保存目录:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        self.result_dir_var = tk.StringVar()
        result_dir_frame = tk.Frame(save_frame, bg=self.bg_color)
        result_dir_frame.pack(fill=tk.X, padx=10, pady=5)
        
        result_entry = tk.Entry(result_dir_frame, textvariable=self.result_dir_var, 
                               font=("宋体", 10))
        result_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))
        
        result_browse_btn = tk.Button(result_dir_frame, text="浏览", command=self.choose_result_dir,
                                     bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=10)
        result_browse_btn.pack(side=tk.RIGHT)
        
        # 词云图保存
        tk.Label(save_frame, text="词云图保存目录:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(20, 5))
        
        self.cloud_dir_var = tk.StringVar()
        cloud_dir_frame = tk.Frame(save_frame, bg=self.bg_color)
        cloud_dir_frame.pack(fill=tk.X, padx=10, pady=5)
        
        cloud_entry = tk.Entry(cloud_dir_frame, textvariable=self.cloud_dir_var, 
                              font=("宋体", 10))
        cloud_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))
        
        cloud_browse_btn = tk.Button(cloud_dir_frame, text="浏览", command=self.choose_cloud_dir,
                                    bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=10)
        cloud_browse_btn.pack(side=tk.RIGHT)
        
        # 保存路径设置按钮
        save_path_btn_frame = tk.Frame(save_frame, bg=self.bg_color)
        save_path_btn_frame.pack(fill=tk.X, padx=10, pady=(20, 10))
        
        save_path_btn = tk.Button(save_path_btn_frame, text="💾 保存路径设置", 
                                 command=self.save_ui_config,
                                 bg="#27ae60", fg="white", font=("微软雅黑", 10), width=20)
        save_path_btn.pack()
    
    def create_vocab_tab(self):
        """创建词汇配置标签页"""
        vocab_tab = tk.Frame(self.notebook, bg=self.bg_color)
        self.notebook.add(vocab_tab, text="🔧 词汇配置")
        
        # 使用两栏布局
        main_frame = tk.Frame(vocab_tab, bg=self.bg_color)
        main_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # 左侧停用词管理
        left_frame = tk.Frame(main_frame, bg=self.bg_color)
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 10))
        
        stopwords_frame = tk.LabelFrame(left_frame, text="停用词管理", 
                                       font=("微软雅黑", 12, "bold"),
                                       bg=self.bg_color, fg=self.text_color)
        stopwords_frame.pack(fill=tk.BOTH, expand=True)
        
        # 停用词列表
        list_frame = tk.Frame(stopwords_frame, bg=self.bg_color)
        list_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.stopwords_listbox = tk.Listbox(list_frame, font=("宋体", 10), selectmode=tk.EXTENDED)
        self.stopwords_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        scrollbar = tk.Scrollbar(list_frame, command=self.stopwords_listbox.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.stopwords_listbox.config(yscrollcommand=scrollbar.set)
        
        # 更新停用词列表
        self.update_stopwords_listbox()
        
        # 停用词按钮
        btn_frame = tk.Frame(stopwords_frame, bg=self.bg_color)
        btn_frame.pack(fill=tk.X, padx=10, pady=10)
        
        tk.Button(btn_frame, text="管理停用词", command=self.manage_stopwords,
                 bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=15).pack()
        
        # 右侧自定义词管理
        right_frame = tk.Frame(main_frame, bg=self.bg_color)
        right_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True, padx=(10, 0))
        
        customwords_frame = tk.LabelFrame(right_frame, text="自定义词管理", 
                                         font=("微软雅黑", 12, "bold"),
                                         bg=self.bg_color, fg=self.text_color)
        customwords_frame.pack(fill=tk.BOTH, expand=True)
        
        # 自定义词列表
        list_frame2 = tk.Frame(customwords_frame, bg=self.bg_color)
        list_frame2.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.customwords_listbox = tk.Listbox(list_frame2, font=("宋体", 10), selectmode=tk.EXTENDED)
        self.customwords_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        scrollbar2 = tk.Scrollbar(list_frame2, command=self.customwords_listbox.yview)
        scrollbar2.pack(side=tk.RIGHT, fill=tk.Y)
        self.customwords_listbox.config(yscrollcommand=scrollbar2.set)
        
        # 更新自定义词列表
        self.update_customwords_listbox()
        
        # 自定义词按钮
        btn_frame2 = tk.Frame(customwords_frame, bg=self.bg_color)
        btn_frame2.pack(fill=tk.X, padx=10, pady=10)
        
        tk.Button(btn_frame2, text="管理自定义词", command=self.manage_customwords,
                 bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=15).pack()
    
    def create_pos_tab(self):
        """创建词性管理标签页"""
        pos_tab = tk.Frame(self.notebook, bg=self.bg_color)
        self.notebook.add(pos_tab, text="🏷️ 词性管理")
        
        # 主框架
        main_frame = tk.Frame(pos_tab, bg=self.bg_color)
        main_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # 词性对照表显示区域
        pos_frame = tk.LabelFrame(main_frame, text="词性对照表管理", 
                                 font=("微软雅黑", 12, "bold"),
                                 bg=self.bg_color, fg=self.text_color)
        pos_frame.pack(fill=tk.BOTH, expand=True)
        
        # 搜索框
        search_frame = tk.Frame(pos_frame, bg=self.bg_color)
        search_frame.pack(fill=tk.X, padx=10, pady=10)
        
        tk.Label(search_frame, text="搜索词性:", bg=self.bg_color, 
                font=("微软雅黑", 11)).pack(side=tk.LEFT)
        self.pos_search_var = tk.StringVar()
        pos_search_entry = tk.Entry(search_frame, textvariable=self.pos_search_var, 
                                   font=("宋体", 11), width=30)
        pos_search_entry.pack(side=tk.LEFT, padx=10)
        
        # 搜索功能
        def on_pos_search(*args):
            keyword = self.pos_search_var.get().strip().lower()
            self.update_pos_listbox(keyword)
        
        # 使用兼容的方式绑定事件
        try:
            self.pos_search_var.trace_add("write", on_pos_search)
        except AttributeError:
            self.pos_search_var.trace("w", on_pos_search)
        
        # 词性列表框
        listbox_frame = tk.Frame(pos_frame, bg=self.bg_color)
        listbox_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # 创建Treeview显示词性对照表
        columns = ("词性标记", "中文含义", "示例")
        self.pos_tree = ttk.Treeview(listbox_frame, columns=columns, show="headings", height=15)
        
        # 设置列标题
        self.pos_tree.heading("词性标记", text="词性标记")
        self.pos_tree.heading("中文含义", text="中文含义")
        self.pos_tree.heading("示例", text="示例")
        
        # 设置列宽
        self.pos_tree.column("词性标记", width=100)
        self.pos_tree.column("中文含义", width=150)
        self.pos_tree.column("示例", width=200)
        
        # 滚动条
        scrollbar = ttk.Scrollbar(listbox_frame, orient="vertical", command=self.pos_tree.yview)
        self.pos_tree.configure(yscrollcommand=scrollbar.set)
        
        # 布局
        self.pos_tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 填充词性对照表
        self.update_pos_listbox("")
        
        # 操作按钮区域
        btn_frame = tk.Frame(pos_frame, bg=self.bg_color)
        btn_frame.pack(fill=tk.X, padx=10, pady=10)
        
        # 添加新词性按钮
        add_pos_btn = tk.Button(btn_frame, text="添加新词性", 
                               command=self.add_new_pos,
                               bg=self.btn_color, fg="white", 
                               font=("微软雅黑", 10), width=15)
        add_pos_btn.pack(side=tk.LEFT, padx=5)
        
        # 编辑选中词性按钮
        edit_pos_btn = tk.Button(btn_frame, text="编辑选中词性", 
                                command=self.edit_selected_pos,
                                bg=self.btn_color, fg="white", 
                                font=("微软雅黑", 10), width=15)
        edit_pos_btn.pack(side=tk.LEFT, padx=5)
        
        # 删除选中词性按钮
        delete_pos_btn = tk.Button(btn_frame, text="删除选中词性", 
                                  command=self.delete_selected_pos,
                                  bg="#e74c3c", fg="white", 
                                  font=("微软雅黑", 10), width=15)
        delete_pos_btn.pack(side=tk.LEFT, padx=5)
        
        # 导入导出按钮区域
        io_frame = tk.Frame(pos_frame, bg=self.bg_color)
        io_frame.pack(fill=tk.X, padx=10, pady=(0, 10))
        
        # 导入词性对照表
        import_pos_btn = tk.Button(io_frame, text="导入词性对照表", 
                                  command=self.import_pos_map,
                                  bg=self.btn_color, fg="white", 
                                  font=("微软雅黑", 10), width=15)
        import_pos_btn.pack(side=tk.LEFT, padx=5)
        
        # 导出词性对照表
        export_pos_btn = tk.Button(io_frame, text="导出词性对照表", 
                                  command=self.export_pos_map,
                                  bg=self.btn_color, fg="white", 
                                  font=("微软雅黑", 10), width=15)
        export_pos_btn.pack(side=tk.LEFT, padx=5)
        
        # 重置为默认
        reset_pos_btn = tk.Button(io_frame, text="重置为默认", 
                                 command=self.reset_pos_map,
                                 bg="#e67e22", fg="white", 
                                 font=("微软雅黑", 10), width=15)
        reset_pos_btn.pack(side=tk.LEFT, padx=5)
        
        # 说明文本
        info_frame = tk.Frame(pos_frame, bg=self.bg_color)
        info_frame.pack(fill=tk.X, padx=10, pady=(0, 10))
        
        info_text = tk.Label(info_frame, 
                            text="提示: 词性标记来自jieba分词，您可以添加、编辑或删除词性映射。\n"
                                 "修改后的词性对照表将保存在用户配置目录中。",
                            bg=self.bg_color, font=("微软雅黑", 9), justify=tk.LEFT)
        info_text.pack()
    
    def create_wordcloud_tab(self):
        """创建词云设置标签页"""
        wordcloud_tab = tk.Frame(self.notebook, bg=self.bg_color)
        self.notebook.add(wordcloud_tab, text="☁️ 词云设置")
        
        # 使用Canvas和Frame实现滚动
        canvas = tk.Canvas(wordcloud_tab, bg=self.bg_color, highlightthickness=0)
        scrollbar = ttk.Scrollbar(wordcloud_tab, orient="vertical", command=canvas.yview)
        scrollable_frame = tk.Frame(canvas, bg=self.bg_color)
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        # 布局
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # 基本设置区域
        basic_frame = tk.LabelFrame(scrollable_frame, text="基本设置", 
                                   font=("微软雅黑", 12, "bold"),
                                   bg=self.bg_color, fg=self.text_color)
        basic_frame.pack(fill=tk.X, padx=20, pady=20)
        
        # 是否生成词云
        self.generate_wordcloud_var = tk.BooleanVar(value=self.wordcloud_settings['generate_wordcloud'])
        generate_check = tk.Checkbutton(basic_frame, text="生成词云图", 
                                       variable=self.generate_wordcloud_var,
                                       bg=self.bg_color, font=("微软雅黑", 11),
                                       command=self.mark_settings_unsaved)
        generate_check.pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        # 是否生成PPT
        self.generate_ppt_var = tk.BooleanVar(value=self.wordcloud_settings['generate_ppt'])
        ppt_check = tk.Checkbutton(basic_frame, text="生成PPT报告", 
                                  variable=self.generate_ppt_var,
                                  bg=self.bg_color, font=("微软雅黑", 11),
                                  command=self.mark_settings_unsaved)
        ppt_check.pack(anchor=tk.W, padx=10, pady=5)
        
        # 输出格式
        tk.Label(basic_frame, text="输出格式:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        self.output_format_var = tk.StringVar(value=self.wordcloud_settings['output_format'])
        format_frame = tk.Frame(basic_frame, bg=self.bg_color)
        format_frame.pack(fill=tk.X, padx=10, pady=5)
        
        format_options = ['png', 'jpg', 'pdf'] if HAS_PDF_SUPPORT else ['png', 'jpg']
        for option in format_options:
            tk.Radiobutton(format_frame, text=option.upper(), 
                          variable=self.output_format_var, value=option,
                          bg=self.bg_color, font=("微软雅黑", 10),
                          command=self.mark_settings_unsaved).pack(side=tk.LEFT, padx=10)
        
        if not HAS_PDF_SUPPORT:
            tk.Label(format_frame, text="(PDF需要安装matplotlib PDF后端)", 
                    bg=self.bg_color, font=("微软雅黑", 9), fg="red").pack(side=tk.LEFT, padx=10)
        
        # 字体设置区域
        font_frame = tk.LabelFrame(scrollable_frame, text="字体设置", 
                                  font=("微软雅黑", 12, "bold"),
                                  bg=self.bg_color, fg=self.text_color)
        font_frame.pack(fill=tk.X, padx=20, pady=(0, 20))
        
        # 字体路径
        tk.Label(font_frame, text="字体文件路径:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        font_path_frame = tk.Frame(font_frame, bg=self.bg_color)
        font_path_frame.pack(fill=tk.X, padx=10, pady=5)
        
        self.font_path_var = tk.StringVar(value=self.wordcloud_settings.get('font_path', ''))
        font_entry = tk.Entry(font_path_frame, textvariable=self.font_path_var, 
                             font=("宋体", 10), width=40)
        font_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))
        
        tk.Button(font_path_frame, text="浏览", command=self.choose_font_file,
                 bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=10).pack(side=tk.RIGHT)
        
        # 字体大小范围
        tk.Label(font_frame, text="字体大小范围:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        font_size_frame = tk.Frame(font_frame, bg=self.bg_color)
        font_size_frame.pack(fill=tk.X, padx=10, pady=5)
        
        # 最大字体
        tk.Label(font_size_frame, text="最大:", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT)
        
        max_font_scale_frame = tk.Frame(font_size_frame, bg=self.bg_color)
        max_font_scale_frame.pack(side=tk.LEFT, padx=5)
        
        self.max_font_size_var = tk.IntVar(value=self.wordcloud_settings['max_font_size'])
        max_font_scale = tk.Scale(max_font_scale_frame, from_=10, to=500, orient=tk.HORIZONTAL,
                                 variable=self.max_font_size_var, bg=self.bg_color,
                                 length=100, command=self.mark_settings_unsaved)
        max_font_scale.pack(side=tk.LEFT)
        
        max_font_entry_frame = tk.Frame(font_size_frame, bg=self.bg_color)
        max_font_entry_frame.pack(side=tk.LEFT, padx=5)
        
        self.max_font_entry = tk.Entry(max_font_entry_frame, 
                                      font=("宋体", 10), width=6, justify='center')
        self.max_font_entry.insert(0, str(self.max_font_size_var.get()))
        self.max_font_entry.pack(side=tk.LEFT)
        tk.Label(max_font_entry_frame, text="px", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT, padx=(5, 20))
        
        # 绑定事件 - 使用失去焦点时验证
        def validate_max_font(event=None):
            try:
                value = int(self.max_font_entry.get())
                if 10 <= value <= 500:
                    self.max_font_size_var.set(value)
                    max_font_scale.set(value)
                else:
                    messagebox.showwarning("警告", "请输入10-500之间的整数")
                    self.max_font_entry.delete(0, tk.END)
                    self.max_font_entry.insert(0, str(self.max_font_size_var.get()))
            except ValueError:
                # 允许空值或正在输入
                pass
        
        self.max_font_entry.bind('<FocusOut>', validate_max_font)
        self.max_font_entry.bind('<Return>', validate_max_font)
        
        # 最小字体
        tk.Label(font_size_frame, text="最小:", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT)
        
        min_font_scale_frame = tk.Frame(font_size_frame, bg=self.bg_color)
        min_font_scale_frame.pack(side=tk.LEFT, padx=5)
        
        self.min_font_size_var = tk.IntVar(value=self.wordcloud_settings['min_font_size'])
        min_font_scale = tk.Scale(min_font_scale_frame, from_=5, to=200, orient=tk.HORIZONTAL,
                                 variable=self.min_font_size_var, bg=self.bg_color,
                                 length=100, command=self.mark_settings_unsaved)
        min_font_scale.pack(side=tk.LEFT)
        
        min_font_entry_frame = tk.Frame(font_size_frame, bg=self.bg_color)
        min_font_entry_frame.pack(side=tk.LEFT, padx=5)
        
        self.min_font_entry = tk.Entry(min_font_entry_frame, 
                                      font=("宋体", 10), width=6, justify='center')
        self.min_font_entry.insert(0, str(self.min_font_size_var.get()))
        self.min_font_entry.pack(side=tk.LEFT)
        tk.Label(min_font_entry_frame, text="px", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT, padx=5)
        
        # 绑定事件
        def validate_min_font(event=None):
            try:
                value = int(self.min_font_entry.get())
                if 5 <= value <= 200:
                    self.min_font_size_var.set(value)
                    min_font_scale.set(value)
                else:
                    messagebox.showwarning("警告", "请输入5-200之间的整数")
                    self.min_font_entry.delete(0, tk.END)
                    self.min_font_entry.insert(0, str(self.min_font_size_var.get()))
            except ValueError:
                pass
        
        self.min_font_entry.bind('<FocusOut>', validate_min_font)
        self.min_font_entry.bind('<Return>', validate_min_font)
        
        # 词数设置 - 改进：更人性化的输入框
        tk.Label(font_frame, text="最大词数:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        max_words_frame = tk.Frame(font_frame, bg=self.bg_color)
        max_words_frame.pack(fill=tk.X, padx=10, pady=5)
        
        # 滑动条
        self.max_words_var = tk.IntVar(value=self.wordcloud_settings['max_words'])
        max_words_scale = tk.Scale(max_words_frame, from_=10, to=500, orient=tk.HORIZONTAL,
                                  variable=self.max_words_var, bg=self.bg_color,
                                  length=250, command=self.mark_settings_unsaved)
        max_words_scale.pack(side=tk.LEFT)
        
        # 输入框
        max_words_entry_frame = tk.Frame(max_words_frame, bg=self.bg_color)
        max_words_entry_frame.pack(side=tk.LEFT, padx=(20, 0))
        
        tk.Label(max_words_entry_frame, text="输入:", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT)
        
        self.max_words_entry = tk.Entry(max_words_entry_frame,
                                       font=("宋体", 10), width=10, justify='center')
        self.max_words_entry.insert(0, str(self.max_words_var.get()))
        self.max_words_entry.pack(side=tk.LEFT, padx=5)
        
        # 绑定事件
        def validate_max_words(event=None):
            try:
                value = int(self.max_words_entry.get())
                if 10 <= value <= 500:
                    self.max_words_var.set(value)
                    max_words_scale.set(value)
                else:
                    messagebox.showwarning("警告", "请输入10-500之间的整数")
                    self.max_words_entry.delete(0, tk.END)
                    self.max_words_entry.insert(0, str(self.max_words_var.get()))
            except ValueError:
                pass
        
        self.max_words_entry.bind('<FocusOut>', validate_max_words)
        self.max_words_entry.bind('<Return>', validate_max_words)
        
        # 颜色设置区域
        color_frame = tk.LabelFrame(scrollable_frame, text="颜色设置", 
                                   font=("微软雅黑", 12, "bold"),
                                   bg=self.bg_color, fg=self.text_color)
        color_frame.pack(fill=tk.X, padx=20, pady=(0, 20))
        
        # 背景颜色 - 添加透明色选项
        tk.Label(color_frame, text="背景颜色:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        self.bg_color_var = tk.StringVar(value=self.wordcloud_settings['background_color'])
        bg_color_frame = tk.Frame(color_frame, bg=self.bg_color)
        bg_color_frame.pack(fill=tk.X, padx=10, pady=5)
        
        bg_colors = ['white', 'black', 'transparent', '#f0f0f0', '#e8f4f8']  # 添加transparent
        for color in bg_colors:
            tk.Radiobutton(bg_color_frame, text=color, 
                          variable=self.bg_color_var, value=color,
                          bg=self.bg_color, font=("微软雅黑", 10),
                          command=self.on_bg_color_change).pack(side=tk.LEFT, padx=10)
        
        # 颜色映射
        tk.Label(color_frame, text="颜色映射:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        self.colormap_var = tk.StringVar(value=self.wordcloud_settings['colormap'])
        colormap_frame = tk.Frame(color_frame, bg=self.bg_color)
        colormap_frame.pack(fill=tk.X, padx=10, pady=5)
        
        colormaps = ['viridis', 'plasma', 'inferno', 'magma', 'cividis', 
                    'spring', 'summer', 'autumn', 'winter', 'cool', 'hot']
        colormap_combo = ttk.Combobox(colormap_frame, textvariable=self.colormap_var,
                                     values=colormaps, state="readonly", 
                                     font=("宋体", 10), width=20)
        colormap_combo.pack(side=tk.LEFT)
        colormap_combo.bind('<<ComboboxSelected>>', lambda e: self.mark_settings_unsaved())
        
        # 图片设置区域
        self.image_frame = tk.LabelFrame(scrollable_frame, text="背景图片设置", 
                                   font=("微软雅黑", 12, "bold"),
                                   bg=self.bg_color, fg=self.text_color)
        self.image_frame.pack(fill=tk.X, padx=20, pady=(0, 20))
        
        # 背景图片
        tk.Label(self.image_frame, text="背景图片:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        image_btn_frame = tk.Frame(self.image_frame, bg=self.bg_color)
        image_btn_frame.pack(fill=tk.X, padx=10, pady=5)
        
        default_bg_btn = tk.Button(image_btn_frame, text="默认背景", 
                                  command=lambda: [self.set_bg_image(None), self.mark_settings_unsaved()],
                                  bg=self.btn_color, fg="white", 
                                  font=("微软雅黑", 10), width=15)
        default_bg_btn.pack(side=tk.LEFT, padx=5)
        
        custom_bg_btn = tk.Button(image_btn_frame, text="选择背景图片", 
                                 command=lambda: [self.choose_bg_image(), self.mark_settings_unsaved()],
                                 bg=self.btn_color, fg="white", 
                                 font=("微软雅黑", 10), width=15)
        custom_bg_btn.pack(side=tk.LEFT, padx=5)
        
        # 显示当前背景图片
        self.bg_image_label = tk.Label(self.image_frame, text="当前: 默认背景", 
                                      bg=self.bg_color, font=("微软雅黑", 10))
        self.bg_image_label.pack(anchor=tk.W, padx=10, pady=5)
        
        # 轮廓设置
        self.image_contour_var = tk.BooleanVar(value=self.wordcloud_settings['image_contour'])
        contour_check = tk.Checkbutton(self.image_frame, text="显示图片轮廓", 
                                      variable=self.image_contour_var,
                                      bg=self.bg_color, font=("微软雅黑", 11),
                                      command=self.on_contour_change)
        contour_check.pack(anchor=tk.W, padx=10, pady=5)
        
        # 轮廓宽度
        tk.Label(self.image_frame, text="轮廓宽度:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(5, 5))
        
        self.contour_width_var = tk.IntVar(value=self.wordcloud_settings['contour_width'])
        contour_width_frame = tk.Frame(self.image_frame, bg=self.bg_color)
        contour_width_frame.pack(fill=tk.X, padx=10, pady=5)
        
        # 滑动条
        contour_width_scale = tk.Scale(contour_width_frame, from_=1, to=10, orient=tk.HORIZONTAL,
                                      variable=self.contour_width_var, bg=self.bg_color,
                                      length=150, command=self.mark_settings_unsaved)
        contour_width_scale.pack(side=tk.LEFT)
        
        # 输入框
        contour_width_entry_frame = tk.Frame(contour_width_frame, bg=self.bg_color)
        contour_width_entry_frame.pack(side=tk.LEFT, padx=(20, 0))
        
        self.contour_width_entry = tk.Entry(contour_width_entry_frame, 
                                           font=("宋体", 10), width=6, justify='center')
        self.contour_width_entry.insert(0, str(self.contour_width_var.get()))
        self.contour_width_entry.pack(side=tk.LEFT)
        tk.Label(contour_width_entry_frame, text="px", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT, padx=5)
        
        # 绑定事件
        def validate_contour_width(event=None):
            try:
                value = int(self.contour_width_entry.get())
                if 1 <= value <= 10:
                    self.contour_width_var.set(value)
                    contour_width_scale.set(value)
                else:
                    messagebox.showwarning("警告", "请输入1-10之间的整数")
                    self.contour_width_entry.delete(0, tk.END)
                    self.contour_width_entry.insert(0, str(self.contour_width_var.get()))
            except ValueError:
                pass
        
        self.contour_width_entry.bind('<FocusOut>', validate_contour_width)
        self.contour_width_entry.bind('<Return>', validate_contour_width)
        
        # 轮廓颜色
        tk.Label(self.image_frame, text="轮廓颜色:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(5, 5))
        
        self.contour_color_var = tk.StringVar(value=self.wordcloud_settings['contour_color'])
        contour_color_frame = tk.Frame(self.image_frame, bg=self.bg_color)
        contour_color_frame.pack(fill=tk.X, padx=10, pady=5)
        
        contour_colors = ['white', 'black', 'red', 'blue', 'green']
        for color in contour_colors:
            tk.Radiobutton(contour_color_frame, text=color, 
                          variable=self.contour_color_var, value=color,
                          bg=self.bg_color, font=("微软雅黑", 10),
                          command=self.mark_settings_unsaved).pack(side=tk.LEFT, padx=10)
        
        # 高级设置区域
        advanced_frame = tk.LabelFrame(scrollable_frame, text="高级设置", 
                                      font=("微软雅黑", 12, "bold"),
                                      bg=self.bg_color, fg=self.text_color)
        advanced_frame.pack(fill=tk.X, padx=20, pady=(0, 20))
        
        # 随机种子
        tk.Label(advanced_frame, text="随机种子:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        self.random_state_var = tk.IntVar(value=self.wordcloud_settings['random_state'])
        random_state_frame = tk.Frame(advanced_frame, bg=self.bg_color)
        random_state_frame.pack(fill=tk.X, padx=10, pady=5)
        
        random_state_entry = tk.Entry(random_state_frame, textvariable=self.random_state_var,
                                     width=10, font=("宋体", 10))
        random_state_entry.pack(side=tk.LEFT)
        random_state_entry.bind('<KeyRelease>', lambda e: self.mark_settings_unsaved())
        
        # 重复词
        self.repeat_var = tk.BooleanVar(value=self.wordcloud_settings['repeat'])
        repeat_check = tk.Checkbutton(advanced_frame, text="允许词重复", 
                                     variable=self.repeat_var,
                                     bg=self.bg_color, font=("微软雅黑", 11),
                                     command=self.mark_settings_unsaved)
        repeat_check.pack(anchor=tk.W, padx=10, pady=5)
        
        # 包含数字
        self.include_numbers_var = tk.BooleanVar(value=self.wordcloud_settings['include_numbers'])
        numbers_check = tk.Checkbutton(advanced_frame, text="包含数字", 
                                      variable=self.include_numbers_var,
                                      bg=self.bg_color, font=("微软雅黑", 11),
                                      command=self.mark_settings_unsaved)
        numbers_check.pack(anchor=tk.W, padx=10, pady=5)
        
        # 水平偏好 - 改进：更人性化的输入框
        tk.Label(advanced_frame, text="水平偏好(0-1):", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=10, pady=(10, 5))
        
        self.prefer_horizontal_var = tk.DoubleVar(value=self.wordcloud_settings['prefer_horizontal'])
        prefer_horizontal_frame = tk.Frame(advanced_frame, bg=self.bg_color)
        prefer_horizontal_frame.pack(fill=tk.X, padx=10, pady=5)
        
        # 滑动条
        prefer_horizontal_scale = tk.Scale(prefer_horizontal_frame, from_=0.0, to=1.0, orient=tk.HORIZONTAL,
                                          variable=self.prefer_horizontal_var, resolution=0.1,
                                          bg=self.bg_color, length=150, command=self.mark_settings_unsaved)
        prefer_horizontal_scale.pack(side=tk.LEFT)
        
        # 输入框
        prefer_horizontal_entry_frame = tk.Frame(prefer_horizontal_frame, bg=self.bg_color)
        prefer_horizontal_entry_frame.pack(side=tk.LEFT, padx=(20, 0))
        
        self.prefer_horizontal_entry = tk.Entry(prefer_horizontal_entry_frame, 
                                               font=("宋体", 10), width=8, justify='center')
        self.prefer_horizontal_entry.insert(0, str(self.prefer_horizontal_var.get()))
        self.prefer_horizontal_entry.pack(side=tk.LEFT)
        
        # 绑定事件
        def validate_prefer_horizontal(event=None):
            try:
                value = float(self.prefer_horizontal_entry.get())
                if 0.0 <= value <= 1.0:
                    self.prefer_horizontal_var.set(value)
                    prefer_horizontal_scale.set(value)
                else:
                    messagebox.showwarning("警告", "请输入0.0-1.0之间的数字")
                    self.prefer_horizontal_entry.delete(0, tk.END)
                    self.prefer_horizontal_entry.insert(0, str(self.prefer_horizontal_var.get()))
            except ValueError:
                pass
        
        self.prefer_horizontal_entry.bind('<FocusOut>', validate_prefer_horizontal)
        self.prefer_horizontal_entry.bind('<Return>', validate_prefer_horizontal)
        
        # 保存设置按钮
        save_btn = tk.Button(scrollable_frame, text="💾 保存词云设置", 
                            command=self.save_wordcloud_settings_with_feedback,
                            bg="#27ae60", fg="white", font=("微软雅黑", 12, "bold"),
                            height=2, width=30)
        save_btn.pack(pady=20)
        
        # 初始化轮廓复选框状态
        self.update_contour_checkbox_state()
    
    def on_bg_color_change(self):
        """背景颜色改变时的处理"""
        self.mark_settings_unsaved()
        
        # 如果是透明背景，禁用轮廓选项
        if self.bg_color_var.get() == 'transparent':
            if self.image_contour_var.get():
                self.image_contour_var.set(False)
                messagebox.showinfo("提示", "透明背景不支持轮廓显示，已自动关闭轮廓选项")
        
        # 更新轮廓复选框状态
        self.update_contour_checkbox_state()
    
    def on_contour_change(self):
        """轮廓设置改变时的处理"""
        self.mark_settings_unsaved()
        
        # 如果选择了透明背景，禁用轮廓
        if self.bg_color_var.get() == 'transparent' and self.image_contour_var.get():
            self.image_contour_var.set(False)
            messagebox.showinfo("提示", "透明背景不支持轮廓显示")
        
        self.update_contour_checkbox_state()
    
    def update_contour_checkbox_state(self):
        """更新轮廓复选框状态"""
        # 如果是透明背景，禁用轮廓复选框
        if self.bg_color_var.get() == 'transparent':
            # 查找轮廓复选框并禁用它
            for widget in self.image_frame.winfo_children():
                if isinstance(widget, tk.Checkbutton):
                    if widget.cget('text') == '显示图片轮廓':
                        widget.config(state='disabled')
                        break
        else:
            # 查找轮廓复选框并启用它
            for widget in self.image_frame.winfo_children():
                if isinstance(widget, tk.Checkbutton):
                    if widget.cget('text') == '显示图片轮廓':
                        widget.config(state='normal')
                        break
    
    def mark_settings_unsaved(self, *args):
        """标记设置未保存"""
        self.settings_saved = False
    
    def save_wordcloud_settings_with_feedback(self):
        """保存词云设置并给予反馈"""
        # 首先验证所有输入框的值
        self.validate_all_entries()
        
        # 更新词云设置
        self.update_wordcloud_settings_from_ui()
        
        # 保存到文件
        if self.save_wordcloud_settings():
            messagebox.showinfo("保存成功", "词云设置已保存！")
            self.settings_saved = True
        else:
            messagebox.showerror("保存失败", "无法保存词云设置")
    
    def validate_all_entries(self):
        """验证所有输入框的值"""
        # 验证最大字体
        try:
            value = int(self.max_font_entry.get())
            if 10 <= value <= 500:
                self.max_font_size_var.set(value)
        except ValueError:
            messagebox.showwarning("警告", "最大字体输入无效，已重置为默认值")
            self.max_font_entry.delete(0, tk.END)
            self.max_font_entry.insert(0, "150")
            self.max_font_size_var.set(150)
        
        # 验证最小字体
        try:
            value = int(self.min_font_entry.get())
            if 5 <= value <= 200:
                self.min_font_size_var.set(value)
        except ValueError:
            messagebox.showwarning("警告", "最小字体输入无效，已重置为默认值")
            self.min_font_entry.delete(0, tk.END)
            self.min_font_entry.insert(0, "10")
            self.min_font_size_var.set(10)
        
        # 验证最大词数
        try:
            value = int(self.max_words_entry.get())
            if 10 <= value <= 500:
                self.max_words_var.set(value)
        except ValueError:
            messagebox.showwarning("警告", "最大词数输入无效，已重置为默认值")
            self.max_words_entry.delete(0, tk.END)
            self.max_words_entry.insert(0, "200")
            self.max_words_var.set(200)
        
        # 验证轮廓宽度
        try:
            value = int(self.contour_width_entry.get())
            if 1 <= value <= 10:
                self.contour_width_var.set(value)
        except ValueError:
            messagebox.showwarning("警告", "轮廓宽度输入无效，已重置为默认值")
            self.contour_width_entry.delete(0, tk.END)
            self.contour_width_entry.insert(0, "1")
            self.contour_width_var.set(1)
        
        # 验证水平偏好
        try:
            value = float(self.prefer_horizontal_entry.get())
            if 0.0 <= value <= 1.0:
                self.prefer_horizontal_var.set(value)
        except ValueError:
            messagebox.showwarning("警告", "水平偏好输入无效，已重置为默认值")
            self.prefer_horizontal_entry.delete(0, tk.END)
            self.prefer_horizontal_entry.insert(0, "0.9")
            self.prefer_horizontal_var.set(0.9)
    
    def save_wordcloud_settings(self):
        """保存词云设置"""
        try:
            # 更新背景图片路径
            self.wordcloud_settings['background_image'] = self.bg_image_path
            
            with open(self.wordcloud_config_file, 'w', encoding='utf-8') as f:
                json.dump(self.wordcloud_settings, f, ensure_ascii=False, indent=2)
            
            self.settings_saved = True
            return True
        except Exception as e:
            print(f"保存词云配置时出错: {e}")
            return False
    
    def choose_font_file(self):
        """选择字体文件"""
        filetypes = [
            ("字体文件", "*.ttf *.ttc *.otf"),
            ("所有文件", "*.*")
        ]
        
        file_path = filedialog.askopenfilename(title="选择字体文件", filetypes=filetypes)
        if file_path:
            self.font_path_var.set(file_path)
            self.mark_settings_unsaved()
    
    def update_pos_listbox(self, keyword):
        """更新词性列表显示"""
        # 清空现有内容
        for item in self.pos_tree.get_children():
            self.pos_tree.delete(item)
        
        # 示例词性
        pos_examples = {
            'n': '苹果、电脑、时间',
            'v': '吃、跑、学习',
            'a': '美丽、快速、聪明',
            'd': '非常、很快、不',
            'm': '一、二、第一',
            'q': '个、只、张',
            'l': '一心一意、马到成功',
            'i': '卧虎藏龙、守株待兔',
            'b': '男、女、主要',
            'vq': '上、下、进、出',
            'ad': '一定、确实、快速'
        }
        
        # 添加词性数据
        for pos_tag, pos_name in sorted(POS_MAP.items()):
            # 搜索过滤
            if keyword and keyword not in pos_tag.lower() and keyword not in pos_name.lower():
                continue
            
            example = pos_examples.get(pos_tag, '')
            self.pos_tree.insert("", tk.END, values=(pos_tag, pos_name, example))
    
    def add_new_pos(self):
        """添加新词性"""
        self.show_pos_editor(None)
    
    def edit_selected_pos(self):
        """编辑选中的词性"""
        selection = self.pos_tree.selection()
        if not selection:
            messagebox.showwarning("警告", "请先选择一个词性")
            return
        
        item = self.pos_tree.item(selection[0])
        pos_tag = item['values'][0]
        self.show_pos_editor(pos_tag)
    
    def delete_selected_pos(self):
        """删除选中的词性"""
        selection = self.pos_tree.selection()
        if not selection:
            messagebox.showwarning("警告", "请先选择一个词性")
            return
        
        item = self.pos_tree.item(selection[0])
        pos_tag = item['values'][0]
        pos_name = item['values'][1]
        
        # 确认删除
        if messagebox.askyesno("确认", f"确定要删除词性 '{pos_tag} ({pos_name})' 吗？"):
            # 只删除用户自定义的词性，不删除内置词性
            if pos_tag in POS_MAP:
                # 从POS_MAP中删除
                del POS_MAP[pos_tag]
                
                # 保存到用户配置
                self.save_pos_map_to_file()
                
                # 更新显示
                self.update_pos_listbox(self.pos_search_var.get())
                
                messagebox.showinfo("成功", f"已删除词性: {pos_tag}")
            else:
                messagebox.showwarning("警告", "无法删除内置词性")
    
    def show_pos_editor(self, pos_tag):
        """显示词性编辑器窗口"""
        editor_window = tk.Toplevel(self.root)
        
        if pos_tag is None:
            editor_window.title("添加新词性")
            is_edit = False
            pos_tag = ""
            pos_name = ""
            example = ""
        else:
            editor_window.title(f"编辑词性: {pos_tag}")
            is_edit = True
            pos_name = POS_MAP.get(pos_tag, "")
            # 示例词性映射
            pos_examples = {
                'n': '苹果、电脑、时间',
                'v': '吃、跑、学习',
                'a': '美丽、快速、聪明',
                'd': '非常、很快、不',
                'm': '一、二、第一',
                'q': '个、只、张'
            }
            example = pos_examples.get(pos_tag, "")
        
        editor_window.geometry("500x350")
        editor_window.configure(bg=self.bg_color)
        editor_window.transient(self.root)
        editor_window.grab_set()
        
        # 词性标记
        tk.Label(editor_window, text="词性标记:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=20, pady=(20, 5))
        
        pos_tag_var = tk.StringVar(value=pos_tag)
        pos_tag_entry = tk.Entry(editor_window, textvariable=pos_tag_var, 
                                font=("宋体", 11), width=30)
        pos_tag_entry.pack(anchor=tk.W, padx=20)
        
        if is_edit:
            pos_tag_entry.config(state='readonly')
        
        # 中文含义
        tk.Label(editor_window, text="中文含义:", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=20, pady=(20, 5))
        
        pos_name_var = tk.StringVar(value=pos_name)
        pos_name_entry = tk.Entry(editor_window, textvariable=pos_name_var, 
                                 font=("宋体", 11), width=30)
        pos_name_entry.pack(anchor=tk.W, padx=20)
        
        # 示例
        tk.Label(editor_window, text="示例(可选):", 
                bg=self.bg_color, font=("微软雅黑", 11)).pack(anchor=tk.W, padx=20, pady=(20, 5))
        
        example_var = tk.StringVar(value=example)
        example_entry = tk.Entry(editor_window, textvariable=example_var, 
                                font=("宋体", 11), width=30)
        example_entry.pack(anchor=tk.W, padx=20)
        
        # 按钮
        button_frame = tk.Frame(editor_window, bg=self.bg_color)
        button_frame.pack(fill=tk.X, padx=20, pady=20)
        
        def save_pos():
            """保存词性"""
            new_tag = pos_tag_var.get().strip()
            new_name = pos_name_var.get().strip()
            
            if not new_tag:
                messagebox.showwarning("警告", "词性标记不能为空")
                return
            
            if not new_name:
                messagebox.showwarning("警告", "中文含义不能为空")
                return
            
            # 保存到POS_MAP
            POS_MAP[new_tag] = new_name
            
            # 保存到文件
            self.save_pos_map_to_file()
            
            # 更新显示
            self.update_pos_listbox(self.pos_search_var.get())
            
            # 关闭窗口
            editor_window.destroy()
            
            messagebox.showinfo("成功", f"已保存词性: {new_tag} -> {new_name}")
        
        save_btn = tk.Button(button_frame, text="保存", 
                            command=save_pos,
                            bg=self.btn_color, fg="white", 
                            font=("微软雅黑", 11), width=15)
        save_btn.pack(side=tk.LEFT, padx=5)
        
        cancel_btn = tk.Button(button_frame, text="取消", 
                              command=editor_window.destroy,
                              bg="#7f8c8d", fg="white", 
                              font=("微软雅黑", 11), width=15)
        cancel_btn.pack(side=tk.LEFT, padx=5)
    
    def save_pos_map_to_file(self):
        """保存词性映射到文件"""
        try:
            with open(self.pos_map_file, 'w', encoding='utf-8') as f:
                json.dump(POS_MAP, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"保存词性映射时出错: {e}")
    
    def import_pos_map(self):
        """导入词性对照表"""
        filetypes = [("JSON文件", "*.json"), ("所有文件", "*.*")]
        file_path = filedialog.askopenfilename(title="选择词性对照表文件", filetypes=filetypes)
        
        if file_path:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    imported_pos_map = json.load(f)
                
                # 更新POS_MAP
                POS_MAP.update(imported_pos_map)
                
                # 保存到用户配置
                self.save_pos_map_to_file()
                
                # 更新显示
                self.update_pos_listbox(self.pos_search_var.get())
                
                messagebox.showinfo("导入成功", f"已导入 {len(imported_pos_map)} 个词性映射")
                
            except Exception as e:
                messagebox.showerror("导入失败", f"无法导入文件: {str(e)}")
    
    def export_pos_map(self):
        """导出词性对照表"""
        filetypes = [("JSON文件", "*.json"), ("所有文件", "*.*")]
        file_path = filedialog.asksaveasfilename(
            title="导出词性对照表",
            defaultextension=".json",
            filetypes=filetypes,
            initialfile="pos_map.json"
        )
        
        if file_path:
            try:
                with open(file_path, 'w', encoding='utf-8') as f:
                    json.dump(POS_MAP, f, ensure_ascii=False, indent=2)
                
                messagebox.showinfo("导出成功", f"已导出 {len(POS_MAP)} 个词性映射到:\n{file_path}")
            except Exception as e:
                messagebox.showerror("导出失败", f"无法导出文件: {str(e)}")
    
    def reset_pos_map(self):
        """重置为默认词性对照表"""
        if messagebox.askyesno("确认", "确定要重置词性对照表吗？这将删除所有自定义词性。"):
            try:
                # 删除用户配置文件
                if os.path.exists(self.pos_map_file):
                    os.remove(self.pos_map_file)
                
                # 重新加载默认配置
                self.load_user_config()
                
                # 更新显示
                self.update_pos_listbox(self.pos_search_var.get())
                
                messagebox.showinfo("重置成功", "词性对照表已重置为默认值")
            except Exception as e:
                messagebox.showerror("重置失败", f"重置时出错: {str(e)}")
    
    def update_stopwords_listbox(self):
        """更新停用词列表框"""
        self.stopwords_listbox.delete(0, tk.END)
        for word in sorted(self.stop_words):
            self.stopwords_listbox.insert(tk.END, word)
        self.stopwords_listbox.insert(0, f"--- 共 {len(self.stop_words)} 个停用词 ---")
    
    def update_customwords_listbox(self):
        """更新自定义词列表框"""
        self.customwords_listbox.delete(0, tk.END)
        for word in sorted(self.custom_words):
            self.customwords_listbox.insert(tk.END, word)
        self.customwords_listbox.insert(0, f"--- 共 {len(self.custom_words)} 个自定义词 ---")
    
    def create_result_tab(self):
        """创建结果显示标签页"""
        result_tab = tk.Frame(self.notebook, bg=self.bg_color)
        self.notebook.add(result_tab, text="📈 分析结果")
        
        # 结果显示区域
        result_frame = tk.LabelFrame(result_tab, text="分析结果", 
                                    font=("微软雅黑", 12, "bold"),
                                    bg=self.bg_color, fg=self.text_color)
        result_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)
        
        # 创建文本区域显示结果
        self.result_text = scrolledtext.ScrolledText(result_frame, 
                                                    font=("宋体", 11),
                                                    wrap=tk.WORD)
        self.result_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # 添加一些初始提示
        self.result_text.insert(tk.END, "欢迎使用词频分析工具 v7.8！\n")
        self.result_text.insert(tk.END, "="*50 + "\n\n")
        self.result_text.insert(tk.END, "使用步骤：\n")
        self.result_text.insert(tk.END, "1. 在'文件管理'标签页添加要分析的文档\n")
        self.result_text.insert(tk.END, "2. 在'分析设置'标签页配置分析参数（包括词性分析）\n")
        self.result_text.insert(tk.END, "3. 在'词汇配置'标签页管理停用词和自定义词\n")
        self.result_text.insert(tk.END, "4. 在'词性管理'标签页可以查看和编辑词性对照表\n")
        self.result_text.insert(tk.END, "5. 在'词云设置'标签页配置词云生成选项\n")
        self.result_text.insert(tk.END, "6. 点击底部的'开始分析'按钮\n\n")
        self.result_text.insert(tk.END, "分析结果将显示词频、词性等信息。\n")
    
    def create_bottom_panel(self):
        """创建底部面板（状态栏和操作按钮）"""
        bottom_frame = tk.Frame(self.root, bg=self.bg_color, height=100)
        bottom_frame.pack(fill=tk.X, side=tk.BOTTOM)
        bottom_frame.pack_propagate(False)
        
        # 进度条
        progress_frame = tk.Frame(bottom_frame, bg=self.bg_color)
        progress_frame.pack(fill=tk.X, padx=20, pady=(10, 5))
        
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(progress_frame, variable=self.progress_var,
                                          maximum=100, mode='determinate', length=800)
        self.progress_bar.pack()
        
        # 按钮和状态栏区域
        control_frame = tk.Frame(bottom_frame, bg=self.bg_color)
        control_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=5)
        
        # 左侧状态栏
        self.status_var = tk.StringVar(value="就绪")
        status_label = tk.Label(control_frame, textvariable=self.status_var,
                               bg=self.bg_color, fg=self.text_color,
                               font=("微软雅黑", 10), anchor=tk.W)
        status_label.pack(side=tk.LEFT, fill=tk.X, expand=True)
        
        # 右侧分析按钮
        self.analyze_btn = tk.Button(control_frame, text="🚀 开始分析", command=self.start_analysis,
                                    bg="#2ecc71", fg="white", font=("微软雅黑", 14, "bold"),
                                    height=2, width=20, relief=tk.RAISED, bd=3)
        self.analyze_btn.pack(side=tk.RIGHT, padx=(0, 10))
        
        # 添加一个配置保存按钮
        save_config_btn = tk.Button(control_frame, text="💾 保存配置", command=self.save_all_config,
                                   bg=self.btn_color, fg="white", font=("微软雅黑", 10),
                                   height=2, width=15)
        save_config_btn.pack(side=tk.RIGHT, padx=10)
    
    def add_files(self):
        """添加文件"""
        filetypes = [
            ("所有支持的文件", "*.txt *.docx *.xlsx *.xls *.pptx"),
            ("文本文件", "*.txt"),
            ("Word文档", "*.docx"),
            ("Excel文件", "*.xlsx *.xls"),
            ("PPT文件", "*.pptx"),
            ("所有文件", "*.*")
        ]
        
        files = filedialog.askopenfilenames(title="选择文档", filetypes=filetypes)
        
        for file in files:
            if file not in self.file_paths:
                self.file_paths.append(file)
                display_name = os.path.basename(file)
                if len(display_name) > 40:
                    display_name = display_name[:37] + "..."
                self.file_listbox.insert(tk.END, f"{display_name}")
        
        self.file_count_label.config(text=f"已选择 {len(self.file_paths)} 个文档")
        self.status_var.set(f"已添加 {len(files)} 个文档")
    
    def remove_file(self):
        """移除选中的文件"""
        selection = self.file_listbox.curselection()
        if selection:
            # 从后往前删除，避免索引变化
            for index in reversed(selection):
                self.file_listbox.delete(index)
                self.file_paths.pop(index)
            
            self.file_count_label.config(text=f"已选择 {len(self.file_paths)} 个文档")
            self.status_var.set(f"已移除 {len(selection)} 个文档")
    
    def choose_result_dir(self):
        """选择结果保存目录"""
        directory = filedialog.askdirectory(title="选择词频结果保存目录")
        if directory:
            self.result_dir_var.set(directory)
            self.status_var.set(f"词频结果保存目录: {directory}")
            # 标记配置需要保存
            self.mark_settings_unsaved()
    
    def choose_cloud_dir(self):
        """选择词云图保存目录"""
        directory = filedialog.askdirectory(title="选择词云图保存目录")
        if directory:
            self.cloud_dir_var.set(directory)
            self.status_var.set(f"词云图保存目录: {directory}")
            # 标记配置需要保存
            self.mark_settings_unsaved()
    
    def set_bg_image(self, image_path=None):
        """设置背景图片"""
        if image_path is None:
            self.bg_image = None
            self.bg_mask = None
            self.bg_image_path = None
            self.bg_image_label.config(text="当前: 默认背景")
            self.status_var.set("已设置为默认背景")
        else:
            try:
                self.bg_image = Image.open(image_path)
                self.bg_mask = np.array(self.bg_image)
                self.bg_image_path = image_path
                filename = os.path.basename(image_path)
                if len(filename) > 20:
                    filename = filename[:17] + "..."
                self.bg_image_label.config(text=f"当前: {filename}")
                self.status_var.set(f"已设置背景图片: {os.path.basename(image_path)}")
            except Exception as e:
                messagebox.showerror("错误", f"无法加载背景图片: {str(e)}")
    
    def choose_bg_image(self):
        """选择背景图片"""
        filetypes = [
            ("图片文件", "*.png *.jpg *.jpeg *.bmp *.gif"),
            ("所有文件", "*.*")
        ]
        
        file_path = filedialog.askopenfilename(title="选择背景图片", filetypes=filetypes)
        if file_path:
            self.set_bg_image(file_path)
    
    def manage_stopwords(self):
        """管理停用词"""
        self.show_word_manager("停用词管理", self.stop_words, "stop_words")
    
    def manage_customwords(self):
        """管理自定义词"""
        self.show_word_manager("自定义词管理", self.custom_words, "custom_words")
    
    def show_word_manager(self, title, word_set, word_type):
        """显示词汇管理窗口"""
        manager_window = tk.Toplevel(self.root)
        manager_window.title(title)
        manager_window.geometry("800x700")
        manager_window.configure(bg=self.bg_color)
        manager_window.transient(self.root)
        manager_window.grab_set()
        
        # 当前词汇显示
        tk.Label(manager_window, text=f"{title}:", 
                bg=self.bg_color, font=("微软雅黑", 14, "bold")).pack(anchor=tk.W, padx=20, pady=(20, 10))
        
        # 添加搜索框
        search_frame = tk.Frame(manager_window, bg=self.bg_color)
        search_frame.pack(fill=tk.X, padx=20, pady=10)
        
        tk.Label(search_frame, text="搜索:", bg=self.bg_color, font=("微软雅黑", 11)).pack(side=tk.LEFT)
        search_var = tk.StringVar()
        search_entry = tk.Entry(search_frame, textvariable=search_var, width=40, font=("宋体", 11))
        search_entry.pack(side=tk.LEFT, padx=10)
        
        # 搜索功能
        def on_search(*args):
            keyword = search_var.get().strip()
            self.update_word_listbox(word_listbox, word_set, keyword)
        
        # 使用兼容的方式绑定事件
        try:
            search_var.trace_add("write", on_search)
        except AttributeError:
            search_var.trace("w", on_search)
        
        # 词汇列表框 - 使用EXTENDED模式支持多选
        listbox_frame = tk.Frame(manager_window, bg=self.bg_color)
        listbox_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        word_listbox = tk.Listbox(listbox_frame, font=("宋体", 11), selectmode=tk.EXTENDED)
        word_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        scrollbar = tk.Scrollbar(listbox_frame, command=word_listbox.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        word_listbox.config(yscrollcommand=scrollbar.set)
        
        # 填充词汇列表
        self.update_word_listbox(word_listbox, word_set, "")
        
        # 操作区域
        operation_frame = tk.Frame(manager_window, bg=self.bg_color)
        operation_frame.pack(fill=tk.X, padx=20, pady=20)
        
        # 添加词汇区域
        add_frame = tk.LabelFrame(operation_frame, text="添加词汇", 
                                 font=("微软雅黑", 11, "bold"),
                                 bg=self.bg_color, fg=self.text_color)
        add_frame.pack(fill=tk.X, pady=(0, 15))
        
        # 单个添加
        single_frame = tk.Frame(add_frame, bg=self.bg_color)
        single_frame.pack(fill=tk.X, padx=10, pady=10)
        
        tk.Label(single_frame, text="单个添加:", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT)
        word_entry = tk.Entry(single_frame, font=("宋体", 10), width=30)
        word_entry.pack(side=tk.LEFT, padx=10)
        
        add_btn = tk.Button(single_frame, text="添加", 
                           command=lambda: self.add_word_to_set(word_entry, word_set, word_listbox, search_var),
                           bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=10)
        add_btn.pack(side=tk.LEFT)
        
        # 批量添加
        batch_frame = tk.Frame(add_frame, bg=self.bg_color)
        batch_frame.pack(fill=tk.X, padx=10, pady=(0, 10))
        
        tk.Label(batch_frame, text="批量添加:", bg=self.bg_color, 
                font=("微软雅黑", 10)).pack(side=tk.LEFT)
        batch_entry = tk.Entry(batch_frame, font=("宋体", 10), width=30)
        batch_entry.pack(side=tk.LEFT, padx=10)
        batch_entry.insert(0, "多个词用逗号或空格分隔")
        
        batch_btn = tk.Button(batch_frame, text="批量添加", 
                             command=lambda: self.add_batch_words(batch_entry, word_set, word_listbox, search_var),
                             bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=10)
        batch_btn.pack(side=tk.LEFT)
        
        # 操作按钮区域
        btn_frame = tk.Frame(operation_frame, bg=self.bg_color)
        btn_frame.pack(fill=tk.X)
        
        # 左侧按钮
        left_frame = tk.Frame(btn_frame, bg=self.bg_color)
        left_frame.pack(side=tk.LEFT, expand=True)
        
        remove_btn = tk.Button(left_frame, text="删除选中", 
                              command=lambda: self.remove_selected_words(word_listbox, word_set, search_var, word_type),
                              bg="#e74c3c", fg="white", font=("微软雅黑", 10), width=12)
        remove_btn.pack(side=tk.LEFT, padx=5)
        
        clear_btn = tk.Button(left_frame, text="清空所有", 
                             command=lambda: self.clear_all_words(word_set, word_listbox, search_var, word_type),
                             bg="#e67e22", fg="white", font=("微软雅黑", 10), width=12)
        clear_btn.pack(side=tk.LEFT, padx=5)
        
        # 右侧按钮
        right_frame = tk.Frame(btn_frame, bg=self.bg_color)
        right_frame.pack(side=tk.RIGHT, expand=True)
        
        import_btn = tk.Button(right_frame, text="导入文件", 
                              command=lambda: self.import_words(word_set, word_listbox, word_type, search_var),
                              bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=12)
        import_btn.pack(side=tk.LEFT, padx=5)
        
        export_btn = tk.Button(right_frame, text="导出文件", 
                              command=lambda: self.export_words(word_set, word_type),
                              bg=self.btn_color, fg="white", font=("微软雅黑", 10), width=12)
        export_btn.pack(side=tk.LEFT, padx=5)
        
        # 保存关闭按钮（单独一行）
        save_frame = tk.Frame(manager_window, bg=self.bg_color)
        save_frame.pack(fill=tk.X, padx=20, pady=(0, 20))
        
        save_btn = tk.Button(save_frame, text="💾 保存并关闭", 
                            command=lambda: self.save_and_close(manager_window, word_type),
                            bg="#27ae60", fg="white", font=("微软雅黑", 12, "bold"),
                            height=2, width=30)
        save_btn.pack()
    
    def add_word_to_set(self, entry_widget, word_set, listbox, search_var):
        """添加单个词到集合"""
        new_word = entry_widget.get().strip()
        if new_word and new_word not in word_set:
            word_set.add(new_word)
            self.update_word_listbox(listbox, word_set, search_var.get())
            entry_widget.delete(0, tk.END)
            if word_set is self.stop_words:
                self.update_stopwords_listbox()
            else:
                self.update_customwords_listbox()
    
    def add_batch_words(self, entry_widget, word_set, listbox, search_var):
        """批量添加词汇"""
        text = entry_widget.get().strip()
        if text:
            import re
            words = re.split(r'[,，\s]+', text)
            added_count = 0
            for word in words:
                word = word.strip()
                if word and word not in word_set and word != "多个词用逗号或空格分隔":
                    word_set.add(word)
                    added_count += 1
            
            if added_count > 0:
                self.update_word_listbox(listbox, word_set, search_var.get())
                entry_widget.delete(0, tk.END)
                if word_set is self.stop_words:
                    self.update_stopwords_listbox()
                else:
                    self.update_customwords_listbox()
                messagebox.showinfo("成功", f"批量添加了 {added_count} 个词汇")
    
    def remove_selected_words(self, listbox, word_set, search_var, word_type):
        """删除选中的词汇（支持多选）"""
        selection = listbox.curselection()
        if not selection:
            messagebox.showwarning("警告", "请先选择要删除的词汇")
            return
        
        # 获取所有选中的词汇，排除标题行
        selected_words = []
        for index in selection:
            word = listbox.get(index)
            if not word.startswith("---"):  # 跳过标题行
                selected_words.append(word)
        
        if not selected_words:
            return
        
        # 确认删除
        delete_count = len(selected_words)
        if delete_count == 1:
            confirm_msg = f"确定要删除词汇 '{selected_words[0]}' 吗？"
        else:
            confirm_msg = f"确定要删除选中的 {delete_count} 个词汇吗？"
        
        if not messagebox.askyesno("确认删除", confirm_msg):
            return
        
        # 从集合中删除选中的词汇
        for word in selected_words:
            if word in word_set:
                word_set.remove(word)
        
        # 更新列表框
        self.update_word_listbox(listbox, word_set, search_var.get())
        
        # 更新主界面列表
        if word_set is self.stop_words:
            self.update_stopwords_listbox()
        else:
            self.update_customwords_listbox()
        
        # 显示成功消息
        messagebox.showinfo("删除成功", f"已删除 {delete_count} 个词汇")
    
    def clear_all_words(self, word_set, listbox, search_var, word_type):
        """清空所有词汇"""
        if not word_set:
            messagebox.showinfo("提示", "词汇列表已经是空的")
            return
        
        if messagebox.askyesno("确认", f"确定要清空所有词汇吗？\n这将删除 {len(word_set)} 个词汇"):
            word_set.clear()
            self.update_word_listbox(listbox, word_set, search_var.get())
            
            # 更新主界面列表
            if word_set is self.stop_words:
                self.update_stopwords_listbox()
            else:
                self.update_customwords_listbox()
            
            messagebox.showinfo("清空成功", "已清空所有词汇")
    
    def update_word_listbox(self, listbox, word_set, keyword):
        """更新词汇列表框"""
        listbox.delete(0, tk.END)
        sorted_words = sorted(word_set)
        
        if keyword:
            filtered_words = [word for word in sorted_words if keyword in word]
        else:
            filtered_words = sorted_words
        
        for word in filtered_words:
            listbox.insert(tk.END, word)
        
        # 显示计数
        if filtered_words:
            listbox.insert(0, f"--- 共 {len(filtered_words)} 个词汇 ---")
        else:
            listbox.insert(0, "--- 暂无词汇 ---")
    
    def import_words(self, word_set, word_listbox, word_type, search_var):
        """从文件导入词汇"""
        filetypes = [("文本文件", "*.txt"), ("所有文件", "*.*")]
        file_path = filedialog.askopenfilename(title="选择词汇文件", filetypes=filetypes)
        
        if file_path:
            try:
                # 尝试检测编码
                encoding = self.detect_encoding(file_path)
                
                imported_words = []
                
                with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                    for line in f:
                        line = line.strip()
                        if line and not line.startswith('#'):
                            import re
                            words_in_line = re.split(r'[,，\s\n]+', line)
                            imported_words.extend(words_in_line)
                
                imported_words = [word.strip() for word in imported_words if word.strip()]
                
                added_count = 0
                for word in imported_words:
                    if word and word not in word_set:
                        word_set.add(word)
                        added_count += 1
                
                self.update_word_listbox(word_listbox, word_set, search_var.get())
                if word_set is self.stop_words:
                    self.update_stopwords_listbox()
                else:
                    self.update_customwords_listbox()
                
                messagebox.showinfo("导入成功", f"成功导入 {added_count} 个词汇")
                
            except Exception as e:
                messagebox.showerror("导入失败", f"无法导入文件: {str(e)}")
    
    def export_words(self, word_set, word_type):
        """导出词汇到文件"""
        filetypes = [("文本文件", "*.txt"), ("所有文件", "*.*")]
        file_path = filedialog.asksaveasfilename(
            title="导出词汇文件",
            defaultextension=".txt",
            filetypes=filetypes,
            initialfile=f"{word_type}.txt"
        )
        
        if file_path:
            try:
                with open(file_path, 'w', encoding='utf-8') as f:
                    for word in sorted(word_set):
                        f.write(word + "\n")
                
                messagebox.showinfo("导出成功", f"已导出 {len(word_set)} 个词汇到:\n{file_path}")
            except Exception as e:
                messagebox.showerror("导出失败", f"无法导出文件: {str(e)}")
    
    def save_and_close(self, window, word_type):
        """保存配置并关闭窗口"""
        self.save_user_config()
        
        if word_type == "custom_words":
            try:
                temp_dict_file = os.path.join(self.config_dir, "custom_words.txt")
                with open(temp_dict_file, 'w', encoding='utf-8') as f:
                    for word in self.custom_words:
                        f.write(f"{word} 100 n\n")
                
                jieba.load_userdict(temp_dict_file)
                self.status_var.set("自定义词库已更新")
            except Exception as e:
                print(f"加载自定义词典时出错: {e}")
        
        window.destroy()
    
    def start_analysis(self):
        """开始分析"""
        if not self.file_paths:
            messagebox.showwarning("警告", "请先选择要分析的文档！")
            return
        
        # 检查保存目录
        if not self.result_dir_var.get():
            messagebox.showwarning("警告", "请选择词频结果保存目录！")
            return
        
        # 检查词云保存目录
        if self.wordcloud_settings['generate_wordcloud'] and not self.cloud_dir_var.get():
            messagebox.showwarning("警告", "请选择词云图保存目录！")
            return
        
        # 检查PDF支持
        if self.output_format_var.get() == 'pdf' and not HAS_PDF_SUPPORT:
            messagebox.showerror("PDF支持错误", 
                                "PDF导出需要matplotlib PDF后端支持。\n"
                                "请安装所需模块或选择PNG/JPG格式。")
            return
        
        # 检查词云设置是否需要保存
        if not self.settings_saved:
            response = messagebox.askyesnocancel("保存设置", 
                "词云设置尚未保存，是否保存并开始分析？\n"
                "点击'是'保存并开始分析，'否'不保存直接开始分析，'取消'返回修改设置。")
            
            if response is None:  # 取消
                return
            elif response:  # 是，保存设置
                # 保存词云设置
                self.save_wordcloud_settings_with_feedback()
                # 继续执行分析
                self.continue_analysis_after_save()
                return
            else:  # 否，不保存直接开始分析
                # 更新词云设置字典，但不保存到文件
                self.update_wordcloud_settings_from_ui()
        
        # 继续执行分析
        self.continue_analysis()
    
    def continue_analysis_after_save(self):
        """保存设置后继续分析"""
        # 等待一小段时间让保存操作完成
        self.root.after(100, self.continue_analysis)
    
    def continue_analysis(self):
        """继续分析流程"""
        # 切换到结果标签页
        self.notebook.select(5)  # 第6个标签页是结果页
        
        # 禁用分析按钮
        self.analyze_btn.config(state='disabled', text="分析中...")
        self.progress_var.set(0)
        self.status_var.set("正在分析文档...")
        
        # 清空结果区域
        self.result_text.delete(1.0, tk.END)
        self.result_text.insert(tk.END, "正在分析文档，请稍候...\n")
        
        # 获取设置
        try:
            top_n = int(self.top_n_entry.get())
        except:
            top_n = 50
        
        # 创建分析任务
        analysis_task = {
            'type': 'analysis',
            'file_paths': self.file_paths.copy(),
            'top_n': top_n,
            'encoding': self.encoding_var.get(),
            'stop_words': list(self.stop_words),
            'custom_words': list(self.custom_words),
            'enable_pos': self.enable_pos_var.get()
        }
        
        # 将任务加入队列
        self.task_queue.put(analysis_task)
    
    def perform_analysis_task(self, task):
        """执行分析任务（在后台线程中）"""
        try:
            # 提取所有文本
            all_text = ""
            total_files = len(task['file_paths'])
            
            for i, file_path in enumerate(task['file_paths']):
                text = self.extract_text(file_path, task['encoding'])
                if text:
                    all_text += text + "\n"
                
                # 更新进度
                progress = (i + 1) / total_files * 50  # 文件读取占50%
                self.root.after(0, self.update_progress, progress)
            
            if not all_text.strip():
                self.result_queue.put({
                    'type': 'error',
                    'data': "无法从文档中提取文本内容"
                })
                return
            
            # 分词和词频统计
            self.root.after(0, self.update_progress, 75)  # 分词处理占25%
            if task['enable_pos']:
                # 使用词性标注分词
                self.word_freq, self.word_pos_freq, self.word_pos_info = self.analyze_text_with_pos(all_text, task['stop_words'], task['custom_words'])
            else:
                # 使用普通分词
                self.word_freq = self.analyze_text(all_text, task['stop_words'], task['custom_words'])
                self.word_pos_freq = {}
                self.word_pos_info = {}
            
            self.root.after(0, self.update_progress, 100)  # 完成
            
            # 将结果加入队列
            self.result_queue.put({
                'type': 'analysis_complete',
                'data': {
                    'word_freq': self.word_freq,
                    'word_pos_freq': self.word_pos_freq,
                    'word_pos_info': self.word_pos_info,
                    'top_n': task['top_n'],
                    'file_count': len(task['file_paths']),
                    'enable_pos': task['enable_pos']
                }
            })
            
        except Exception as e:
            self.result_queue.put({
                'type': 'error',
                'data': f"分析过程中出错: {str(e)}\n\n错误类型: {type(e).__name__}"
            })
    
    def update_progress(self, value):
        """更新进度条"""
        self.progress_var.set(value)
    
    def extract_text(self, file_path, encoding):
        """从不同格式的文件中提取文本"""
        text = ""
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == '.txt':
                if encoding == "auto":
                    file_encoding = self.detect_encoding(file_path)
                else:
                    file_encoding = encoding
                
                try:
                    with open(file_path, 'r', encoding=file_encoding, errors='ignore') as f:
                        text = f.read()
                except UnicodeDecodeError:
                    # 尝试其他编码
                    for enc in ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin-1']:
                        try:
                            with open(file_path, 'r', encoding=enc, errors='ignore') as f:
                                text = f.read()
                            break
                        except:
                            continue
            
            elif ext == '.docx':
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
                # 提取表格内容
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text += cell.text + " "
            
            elif ext in ['.xlsx', '.xls']:
                try:
                    # 读取所有工作表
                    excel_file = pd.ExcelFile(file_path)
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(excel_file, sheet_name=sheet_name, header=None)
                        for col in df.columns:
                            text += " ".join(df[col].dropna().astype(str).tolist()) + " "
                except Exception as e:
                    print(f"读取Excel文件时出错: {e}")
            
            elif ext == '.pptx':
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            
        except Exception as e:
            print(f"提取文件 {file_path} 时出错: {e}")
        
        return text
    
    def analyze_text(self, text, stop_words, custom_words):
        """分析文本并统计词频（不带词性）"""
        # 使用结巴分词
        words = jieba.lcut(text)
        
        # 创建本地集合
        stop_words_set = set(stop_words)
        custom_words_set = set(custom_words)
        
        # 过滤停用词和非中文字符
        filtered_words = []
        for word in words:
            word = word.strip()
            if word and word not in stop_words_set:
                if (len(word) > 1 or word in custom_words_set) and not word.isdigit():
                    if re.search(r'[\u4e00-\u9fff]', word) or re.search(r'[a-zA-Z]', word):
                        filtered_words.append(word)
        
        # 统计词频
        word_counter = Counter(filtered_words)
        
        return dict(word_counter)
    
    def analyze_text_with_pos(self, text, stop_words, custom_words):
        """分析文本并统计词频（带词性）"""
        # 使用结巴分词进行词性标注
        words_pos = pseg.lcut(text)
        
        # 创建本地集合
        stop_words_set = set(stop_words)
        custom_words_set = set(custom_words)
        
        # 存储词频和词性信息
        word_freq = {}  # 词频统计
        word_pos_counter = Counter()  # 词性频率统计
        word_pos_info = {}  # 每个词的词性信息
        word_pos_temp = {}  # 临时存储词的词性出现次数
        
        for word, pos in words_pos:
            word = word.strip()
            if word and word not in stop_words_set:
                if (len(word) > 1 or word in custom_words_set) and not word.isdigit():
                    if re.search(r'[\u4e00-\u9fff]', word) or re.search(r'[a-zA-Z]', word):
                        # 统计词频
                        word_freq[word] = word_freq.get(word, 0) + 1
                        
                        # 统计词性频率
                        pos_cn = POS_MAP.get(pos, pos)
                        word_pos_counter[pos_cn] += 1
                        
                        # 记录每个词的词性出现情况
                        if word not in word_pos_temp:
                            word_pos_temp[word] = Counter()
                        word_pos_temp[word][pos_cn] += 1
        
        # 确定每个词的最终词性（选择出现次数最多的词性）
        for word, pos_counter in word_pos_temp.items():
            most_common_pos = pos_counter.most_common(1)
            if most_common_pos:
                word_pos_info[word] = most_common_pos[0][0]
        
        return word_freq, dict(word_pos_counter), word_pos_info
    
    def display_results(self, data):
        """在文本框中显示结果"""
        self.result_text.delete(1.0, tk.END)
        
        word_freq = data['word_freq']
        word_pos_freq = data.get('word_pos_freq', {})
        word_pos_info = data.get('word_pos_info', {})
        top_n = data['top_n']
        file_count = data['file_count']
        enable_pos = data.get('enable_pos', False)
        
        # 按词频排序
        sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
        
        # 显示前N个
        display_words = sorted_words[:top_n]
        
        if enable_pos:
            self.result_text.insert(tk.END, f"📊 词频分析结果 (前{top_n}个，带词性分析)\n")
            self.result_text.insert(tk.END, "="*70 + "\n\n")
            
            # 显示词频统计（带词性分析）
            for i, (word, freq) in enumerate(display_words, 1):
                # 获取词的词性
                pos = word_pos_info.get(word, '未知')
                self.result_text.insert(tk.END, f"{i:3d}. {word:15s} ({pos:5s}) : {freq:6d} 次\n")
        else:
            self.result_text.insert(tk.END, f"📊 词频分析结果 (前{top_n}个)\n")
            self.result_text.insert(tk.END, "="*60 + "\n\n")
            
            for i, (word, freq) in enumerate(display_words, 1):
                self.result_text.insert(tk.END, f"{i:3d}. {word:25s} : {freq:6d} 次\n")
        
        # 显示词性分布（如果启用了词性分析）
        if enable_pos and word_pos_freq:
            self.result_text.insert(tk.END, "\n" + "="*60 + "\n")
            self.result_text.insert(tk.END, "📈 词性分布统计\n")
            self.result_text.insert(tk.END, "-"*60 + "\n")
            
            # 按词性频率排序
            sorted_pos = sorted(word_pos_freq.items(), key=lambda x: x[1], reverse=True)
            for pos, freq in sorted_pos[:20]:  # 显示前20个词性
                self.result_text.insert(tk.END, f"{pos:15s} : {freq:6d} 次\n")
        
        # 显示统计信息
        self.result_text.insert(tk.END, "\n" + "="*60 + "\n")
        self.result_text.insert(tk.END, f"总词汇数: {len(word_freq)}\n")
        self.result_text.insert(tk.END, f"总词频数: {sum(word_freq.values())}\n")
        if enable_pos:
            self.result_text.insert(tk.END, f"词性种类: {len(word_pos_freq)}\n")
        self.result_text.insert(tk.END, f"分析文档: {file_count} 个\n")
        self.result_text.insert(tk.END, f"分析时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n")
        
        # 重新启用分析按钮
        self.analyze_btn.config(state='normal', text="🚀 开始分析")
    
    def save_results(self, top_n):
        """保存词频结果 - 修复CSV编码问题"""
        try:
            save_dir = self.result_dir_var.get()
            os.makedirs(save_dir, exist_ok=True)
            
            # 排序词频
            sorted_words = sorted(self.word_freq.items(), key=lambda x: x[1], reverse=True)
            
            # 准备数据框
            word_data = []
            for word, freq in sorted_words:
                # 获取词性信息
                if self.enable_pos_var.get():
                    pos = self.word_pos_info.get(word, '未知')
                else:
                    pos = ''
                
                word_data.append({
                    '词语': word,
                    '词性': pos,
                    '频次': freq
                })
            
            # 保存为CSV - 包含词性列，使用用户选择的编码
            csv_path = os.path.join(save_dir, "word_frequency.csv")
            df = pd.DataFrame(word_data)
            
            # 获取用户选择的CSV编码
            csv_encoding = self.csv_encoding_var.get()
            
            try:
                df.to_csv(csv_path, index=False, encoding=csv_encoding)
                csv_status = f"编码: {csv_encoding}"
            except Exception as e:
                # 如果选择的编码失败，尝试其他编码
                try:
                    df.to_csv(csv_path, index=False, encoding='utf-8-sig')
                    csv_status = f"编码: utf-8-sig (自动切换)"
                except Exception as e2:
                    try:
                        df.to_csv(csv_path, index=False, encoding='gbk')
                        csv_status = f"编码: gbk (自动切换)"
                    except Exception as e3:
                        df.to_csv(csv_path, index=False, encoding='utf-8')
                        csv_status = f"编码: utf-8 (自动切换)"
            
            # 保存词性分布为CSV（如果启用了词性分析）
            if self.enable_pos_var.get() and self.word_pos_freq:
                pos_csv_path = os.path.join(save_dir, "part_of_speech.csv")
                sorted_pos = sorted(self.word_pos_freq.items(), key=lambda x: x[1], reverse=True)
                pos_df = pd.DataFrame(sorted_pos, columns=["词性", "频次"])
                
                try:
                    pos_df.to_csv(pos_csv_path, index=False, encoding=csv_encoding)
                except Exception as e:
                    try:
                        pos_df.to_csv(pos_csv_path, index=False, encoding='utf-8-sig')
                    except Exception as e2:
                        try:
                            pos_df.to_csv(pos_csv_path, index=False, encoding='gbk')
                        except Exception as e3:
                            pos_df.to_csv(pos_csv_path, index=False, encoding='utf-8')
        
            # 保存为TXT（带词性信息）
            txt_path = os.path.join(save_dir, "word_frequency.txt")
            with open(txt_path, 'w', encoding='utf-8') as f:
                f.write("词频分析结果\n")
                f.write("="*70 + "\n\n")
                
                for i, item in enumerate(word_data[:top_n], 1):
                    word = item['词语']
                    freq = item['频次']
                    pos = item['词性']
                    
                    if pos:  # 如果有词性信息
                        f.write(f"{i:3d}. {word:15s} ({pos:5s}) : {freq:6d} 次\n")
                    else:
                        f.write(f"{i:3d}. {word:25s} : {freq:6d} 次\n")
                
                f.write("\n" + "="*60 + "\n")
                f.write(f"总词汇数: {len(self.word_freq)}\n")
                f.write(f"总词频数: {sum(self.word_freq.values())}\n")
                
                # 如果启用了词性分析，保存词性分布
                if self.enable_pos_var.get() and self.word_pos_freq:
                    f.write(f"词性种类: {len(self.word_pos_freq)}\n")
                    f.write("\n词性分布:\n")
                    sorted_pos = sorted(self.word_pos_freq.items(), key=lambda x: x[1], reverse=True)
                    for pos, freq in sorted_pos[:20]:
                        f.write(f"{pos:15s} : {freq:6d} 次\n")
                
                f.write(f"分析时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n")
            
            # 显示保存路径和编码信息
            self.result_text.insert(tk.END, f"\n💾 结果已保存至:\n")
            self.result_text.insert(tk.END, f"CSV文件: {csv_path} ({csv_status})\n")
            self.result_text.insert(tk.END, f"TXT文件: {txt_path}\n")
            
            if self.enable_pos_var.get() and self.word_pos_freq:
                self.result_text.insert(tk.END, f"词性分布CSV: {pos_csv_path}\n")
                
            # 提供编码使用建议
            if csv_encoding != 'utf-8-sig':
                self.result_text.insert(tk.END, f"\n📝 编码使用建议:\n")
                self.result_text.insert(tk.END, "1. utf-8-sig: Excel可以正确识别，推荐使用\n")
                self.result_text.insert(tk.END, "2. gbk: Windows中文系统默认编码\n")
                self.result_text.insert(tk.END, "3. 如果Excel打开乱码，请尝试不同编码\n")
        
        except Exception as e:
            messagebox.showwarning("保存失败", f"无法保存结果: {str(e)}\n\n请检查文件权限和磁盘空间。")
    
    def generate_wordcloud_task(self, task):
        """生成词云图任务（在后台线程中）"""
        try:
            word_freq = task['word_freq']
            top_n = task['top_n']
            settings = task['settings']
            
            # 获取保存目录
            save_dir = self.cloud_dir_var.get()
            os.makedirs(save_dir, exist_ok=True)
            
            # 限制词数
            if top_n < len(word_freq):
                sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
                word_freq = dict(sorted_words[:top_n])
            
            # 获取字体路径
            font_path = settings.get('font_path')
            if not font_path or not os.path.exists(font_path):
                font_path = self.get_font_path()
            
            # 处理透明背景
            background_color = settings['background_color']
            if background_color == 'transparent':
                background_color = None  # wordcloud将使用透明背景
            
            # 如果有背景图片，需要处理图片通道问题
            mask = None
            if self.bg_mask is not None:
                # 修复通道数不匹配的问题
                mask = self.fix_image_channels(self.bg_mask, background_color)
                print(f"图片处理: 原始形状 {self.bg_mask.shape}, 处理后形状 {mask.shape}")
            
            # 创建词云配置
            wc_config = {
                'font_path': font_path,
                'background_color': background_color,
                'max_words': settings['max_words'],
                'max_font_size': settings['max_font_size'],
                'min_font_size': settings['min_font_size'],
                'width': settings['width'],
                'height': settings['height'],
                'random_state': settings['random_state'],
                'repeat': settings['repeat'],
                'include_numbers': settings['include_numbers'],
                'prefer_horizontal': settings['prefer_horizontal'],
                'scale': settings['scale'],
                'colormap': settings['colormap'],
                'contour_width': 0,  # 默认不显示轮廓
                'contour_color': settings['contour_color'],
                'mode': 'RGBA' if background_color is None else 'RGB'
            }
            
            # 如果有背景图片
            if mask is not None:
                wc_config['mask'] = mask
                
                # 如果是透明背景，禁用轮廓
                if background_color is None and settings['image_contour']:
                    print("警告: 透明背景不支持轮廓显示，已禁用轮廓")
                    settings['image_contour'] = False
                
                # 修复轮廓线曲里拐弯的问题 - 使用更高质量的轮廓
                if settings['image_contour']:
                    wc_config['contour_width'] = settings['contour_width']
                    wc_config['contour_color'] = settings['contour_color']
                    # 添加更多优化参数
                    wc_config['margin'] = 2  # 增加边距
                    wc_config['scale'] = 2   # 提高缩放质量
            
            # 创建词云
            wc = WordCloud(**wc_config)
            wc.generate_from_frequencies(word_freq)
            
            # 如果使用背景图片且需要重新着色
            if mask is not None and self.bg_image is not None and background_color is not None:
                try:
                    # 确保颜色生成器的形状与mask一致
                    image_colors = ImageColorGenerator(mask)
                    wc.recolor(color_func=image_colors)
                except Exception as e:
                    print(f"重新着色时出错: {e}")
                    # 如果重新着色失败，继续使用默认颜色
            
            # 保存词云图
            timestamp = int(time.time())
            output_format = settings.get('output_format', 'png')
            
            # 始终生成PNG预览文件用于显示
            png_preview_path = os.path.join(save_dir, f"wordcloud_{timestamp}_preview.png")
            
            # 保存PNG预览文件
            plt.figure(figsize=(12, 8), dpi=150)
            plt.imshow(wc, interpolation="bilinear")
            plt.axis("off")
            plt.tight_layout(pad=0)
            plt.savefig(png_preview_path, dpi=150, bbox_inches='tight', pad_inches=0, format='png')
            plt.close('all')
            
            # 保存主要输出格式的文件
            if output_format == 'png':
                # 对于PNG格式，PNG预览文件就是主文件
                main_save_path = png_preview_path
            elif output_format in ['jpg', 'jpeg']:
                # 保存JPG格式的主文件
                main_save_path = os.path.join(save_dir, f"wordcloud_{timestamp}.jpg")
                plt.figure(figsize=(12, 8), dpi=300)
                plt.imshow(wc, interpolation="bilinear")
                plt.axis("off")
                plt.tight_layout(pad=0)
                plt.savefig(main_save_path, dpi=300, bbox_inches='tight', pad_inches=0, format='jpg')
                plt.close('all')
            elif output_format == 'pdf':
                # 保存PDF格式的主文件
                main_save_path = os.path.join(save_dir, f"wordcloud_{timestamp}.pdf")
                plt.figure(figsize=(12, 8), dpi=300)
                plt.imshow(wc, interpolation="bilinear")
                plt.axis("off")
                plt.tight_layout(pad=0)
                plt.savefig(main_save_path, format='pdf', bbox_inches='tight', pad_inches=0)
                plt.close('all')
            else:
                # 默认保存为PNG
                main_save_path = png_preview_path
            
            result_data = {
                'image_path': main_save_path,
                'png_preview_path': png_preview_path  # 新增：PNG预览文件路径
            }
            
            # 生成PPT报告（如果启用）
            if settings.get('generate_ppt'):
                # 使用词云图保存目录
                ppt_path = self.generate_ppt_report(save_dir, timestamp, word_freq, wc)
                if ppt_path:
                    result_data['ppt_path'] = ppt_path
                    # 也在词频结果目录保存一份副本
                    try:
                        result_dir = self.result_dir_var.get()
                        if result_dir:
                            ppt_copy_path = os.path.join(result_dir, f"wordcloud_report_{timestamp}.pptx")
                            import shutil
                            shutil.copy2(ppt_path, ppt_copy_path)
                    except Exception as e:
                        print(f"复制PPT到结果目录时出错: {e}")
            
            # 将结果加入队列
            self.result_queue.put({
                'type': 'wordcloud_complete',
                'data': result_data
            })
            
        except Exception as e:
            print(f"生成词云图时出错: {e}")
            import traceback
            traceback.print_exc()
            self.result_queue.put({
                'type': 'error',
                'data': f"无法生成词云图: {str(e)}"
            })
    
    def fix_image_channels(self, mask, background_color):
        """修复图片通道数不匹配问题"""
        # 获取图片形状
        height, width = mask.shape[:2]
        
        # 如果是透明背景，需要RGBA模式
        if background_color is None:
            # 如果是4通道，直接返回
            if mask.shape[2] == 4:
                return mask
            # 如果是3通道，转换为4通道
            elif mask.shape[2] == 3:
                # 创建alpha通道（全不透明）
                alpha_channel = np.ones((height, width), dtype=mask.dtype) * 255
                # 合并为4通道
                return np.dstack((mask, alpha_channel))
        else:
            # 如果是非透明背景，需要RGB模式
            # 如果是4通道，转换为3通道
            if mask.shape[2] == 4:
                # 丢弃alpha通道
                return mask[:, :, :3]
            # 如果是3通道，直接返回
            elif mask.shape[2] == 3:
                return mask
        
        # 其他情况，返回原始mask
        return mask
    
    def generate_ppt_report(self, save_dir, timestamp, word_freq, wordcloud):
        """生成PPT报告 - 增强版：添加词性分析页面和高频词汇词性信息"""
        try:
            # 创建PPT
            prs = Presentation()
            
            # 第一页：标题页
            slide_layout = prs.slide_layouts[0]  # 标题幻灯片
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "词频分析报告"
            
            subtitle = slide.placeholders[1]
            enable_pos = self.enable_pos_var.get()
            if enable_pos:
                pos_info = "已启用词性分析"
            else:
                pos_info = "未启用词性分析"
            
            subtitle.text = f"生成时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n总词汇数: {len(word_freq)}\n总频次: {sum(word_freq.values())}\n{pos_info}"
            
            # 第二页：词云图
            slide_layout = prs.slide_layouts[5]  # 只有标题的幻灯片
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "词云图"
            
            # 保存词云图为临时图片用于插入PPT
            temp_img_path = os.path.join(save_dir, f"temp_wordcloud_{timestamp}.png")
            
            # 生成高质量的预览图片
            plt.figure(figsize=(10, 6), dpi=150)
            plt.imshow(wordcloud, interpolation="bilinear")
            plt.axis("off")
            plt.tight_layout(pad=0)
            plt.savefig(temp_img_path, dpi=150, bbox_inches='tight', pad_inches=0)
            plt.close('all')
            
            # 插入图片
            left = Inches(1)
            top = Inches(2)
            height = Inches(5)
            
            try:
                # 确保图片文件存在
                if os.path.exists(temp_img_path) and os.path.getsize(temp_img_path) > 0:
                    pic = slide.shapes.add_picture(temp_img_path, left, top, height=height)
                else:
                    # 如果图片生成失败，创建文本占位符
                    textbox = slide.shapes.add_textbox(left, top, Inches(8), height)
                    tf = textbox.text_frame
                    tf.clear()
                    p = tf.add_paragraph()
                    p.text = "词云图已生成，但无法插入到PPT中。\n请查看保存的图片文件。"
                    p.font.size = Pt(24)
                    p.font.color.rgb = RGBColor(255, 0, 0)
                    p.alignment = PP_ALIGN.CENTER
            except Exception as e:
                print(f"插入图片到PPT时出错: {e}")
                # 创建文本占位符
                textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
                tf = textbox.text_frame
                tf.clear()
                p = tf.add_paragraph()
                p.text = f"词云图已保存为PNG文件\n错误: {str(e)}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(0, 0, 255)
                p.alignment = PP_ALIGN.CENTER
            
            # 第三页：高频词汇（包含词性信息）
            slide_layout = prs.slide_layouts[1]  # 标题和内容的幻灯片
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            if enable_pos:
                title.text = "高频词汇（带词性）"
            else:
                title.text = "高频词汇"
            
            content = slide.shapes.placeholders[1]
            tf = content.text_frame
            tf.clear()  # 清空原有内容
            
            # 添加前20个高频词
            sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
            
            # 如果启用了词性分析，添加词性信息
            if enable_pos:
                for i, (word, freq) in enumerate(sorted_words[:20], 1):
                    p = tf.add_paragraph()
                    pos = self.word_pos_info.get(word, '未知')
                    p.text = f"{i}. {word} ({pos}): {freq}次"
                    p.level = 0
            else:
                for i, (word, freq) in enumerate(sorted_words[:20], 1):
                    p = tf.add_paragraph()
                    p.text = f"{i}. {word}: {freq}次"
                    p.level = 0
            
            # 第四页：词性分布分析（仅当启用了词性分析时）
            if enable_pos and self.word_pos_freq:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                title.text = "词性分布分析"
                
                content = slide.shapes.placeholders[1]
                tf = content.text_frame
                tf.clear()
                
                # 添加词性分布统计
                sorted_pos = sorted(self.word_pos_freq.items(), key=lambda x: x[1], reverse=True)
                total_pos_freq = sum(self.word_pos_freq.values())
                
                p = tf.add_paragraph()
                p.text = f"总词性数: {len(self.word_pos_freq)}"
                p.level = 0
                
                p = tf.add_paragraph()
                p.text = f"总词性频次: {total_pos_freq}"
                p.level = 0
                
                p = tf.add_paragraph()
                p.text = "前15个高频词性:"
                p.level = 0
                
                for pos, freq in sorted_pos[:15]:
                    p = tf.add_paragraph()
                    percentage = (freq / total_pos_freq) * 100
                    p.text = f"• {pos}: {freq}次 ({percentage:.1f}%)"
                    p.level = 1
            
            # 第五页：统计分析（如果启用了词性分析，调整页码）
            analysis_slide_num = 5 if (enable_pos and self.word_pos_freq) else 4
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "统计分析"
            
            content = slide.shapes.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # 添加统计信息
            sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
            stats_text = f"""
            文档数量: {len(self.file_paths)}
            总词汇数: {len(word_freq)}
            总频次数: {sum(word_freq.values())}
            平均频次: {sum(word_freq.values()) / len(word_freq):.2f}
            最高频次: {sorted_words[0][1] if sorted_words else 0} ({sorted_words[0][0] if sorted_words else '无'})
            最低频次: {sorted_words[-1][1] if sorted_words else 0} ({sorted_words[-1][0] if sorted_words else '无'})
            
            词云设置:
            - 最大词数: {self.wordcloud_settings['max_words']}
            - 字体大小: {self.wordcloud_settings['min_font_size']}-{self.wordcloud_settings['max_font_size']}px
            - 背景颜色: {self.wordcloud_settings['background_color']}
            - 颜色映射: {self.wordcloud_settings['colormap']}
            - 输出格式: {self.wordcloud_settings['output_format']}
            """
            
            # 如果启用了词性分析，添加词性统计信息
            if enable_pos and self.word_pos_freq:
                total_pos_freq = sum(self.word_pos_freq.values())
                sorted_pos = sorted(self.word_pos_freq.items(), key=lambda x: x[1], reverse=True)
                
                stats_text += f"""
                
                词性分析统计:
                - 词性种类: {len(self.word_pos_freq)}
                - 高频词性: {sorted_pos[0][0] if sorted_pos else '无'} ({sorted_pos[0][1] if sorted_pos else 0}次)
                - 低频词性: {sorted_pos[-1][0] if sorted_pos else '无'} ({sorted_pos[-1][1] if sorted_pos else 0}次)
                - 平均词性频次: {total_pos_freq / len(self.word_pos_freq):.2f}
                """
            
            for line in stats_text.strip().split('\n'):
                p = tf.add_paragraph()
                p.text = line.strip()
                p.level = 0
            
            # 第六页：词性示例和说明（仅当启用了词性分析时）
            if enable_pos:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                title.text = "词性示例和说明"
                
                content = slide.shapes.placeholders[1]
                tf = content.text_frame
                tf.clear()
                
                # 添加常见词性说明
                pos_examples = {
                    'n': '名词 (苹果、电脑、时间)',
                    'v': '动词 (吃、跑、学习)',
                    'a': '形容词 (美丽、快速、聪明)',
                    'd': '副词 (非常、很快、不)',
                    'm': '数词 (一、二、第一)',
                    'q': '量词 (个、只、张)',
                    'r': '代词 (我、你、他)',
                    'p': '介词 (在、从、向)',
                    'c': '连词 (和、或、但是)',
                    'u': '助词 (的、地、得)',
                    'nr': '人名 (张三、李四、王五)',
                    'ns': '地名 (北京、上海、广州)',
                    'nt': '机构名 (公司、学校、医院)'
                }
                
                p = tf.add_paragraph()
                p.text = "常见词性标记说明:"
                p.level = 0
                
                # 分两列显示词性说明
                left_col = ""
                right_col = ""
                pos_items = list(pos_examples.items())
                mid_point = (len(pos_items) + 1) // 2
                
                for i, (pos_tag, example) in enumerate(pos_items):
                    line = f"{pos_tag}: {example}\n"
                    if i < mid_point:
                        left_col += line
                    else:
                        right_col += line
                
                # 创建两列显示
                cols_frame = content.text_frame
                p_left = cols_frame.add_paragraph()
                p_left.text = left_col
                p_left.level = 1
                
                # 添加一些间距
                p_space = cols_frame.add_paragraph()
                p_space.text = ""
                
                # 如果右侧有内容，添加右侧列
                if right_col:
                    p_right = cols_frame.add_paragraph()
                    p_right.text = right_col
                    p_right.level = 1
            
            # 最后一页：保存信息（根据是否有词性页面调整页码）
            save_slide_num = 7 if (enable_pos and self.word_pos_freq) else 6
            if enable_pos:
                save_slide_num += 1  # 词性示例页
            
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "文件保存信息"
            
            content = slide.shapes.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            # 添加文件保存信息
            wordcloud_filename = f"wordcloud_{timestamp}.{self.wordcloud_settings['output_format']}"
            wordcloud_path = os.path.join(save_dir, wordcloud_filename)
            ppt_filename = f"wordcloud_report_{timestamp}.pptx"
            ppt_path = os.path.join(save_dir, ppt_filename)
            
            save_info = f"""
            词云图已保存至:
            {wordcloud_path}
            
            PPT报告已保存至:
            {ppt_path}
            
            词频分析结果已保存至:
            {self.result_dir_var.get() if self.result_dir_var.get() else '未设置'}
            
            报告生成时间: {time.strftime('%Y-%m-%d %H:%M:%S')}
            
            报告包含页面:
            1. 标题页
            2. 词云图
            3. 高频词汇{'（带词性）' if enable_pos else ''}
            {('4. 词性分布分析' if self.word_pos_freq else '') + ('\n5. 统计分析' if self.word_pos_freq else '4. 统计分析')}
            {('6. 词性示例和说明' if enable_pos else '')}
            {('7. 文件保存信息' if (enable_pos and self.word_pos_freq) else ('6. 文件保存信息' if enable_pos else '5. 文件保存信息'))}
            
            备注:
            1. 词云图支持PNG、JPG、PDF格式
            2. 词频结果包含CSV、TXT格式
            3. 详细统计信息请查看相关文件
            """
            
            # 清理空行
            save_info_lines = [line for line in save_info.strip().split('\n') if line.strip()]
            save_info_clean = '\n'.join(save_info_lines)
            
            for line in save_info_clean.strip().split('\n'):
                p = tf.add_paragraph()
                p.text = line.strip()
                p.level = 0
            
            # 保存PPT
            ppt_path = os.path.join(save_dir, ppt_filename)
            prs.save(ppt_path)
            
            # 删除临时图片
            try:
                if os.path.exists(temp_img_path):
                    os.remove(temp_img_path)
            except:
                pass
            
            print(f"PPT报告已生成: {ppt_path}")
            return ppt_path
            
        except Exception as e:
            print(f"生成PPT报告时出错: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def get_font_path(self):
        """获取中文字体路径"""
        # 常见字体路径
        font_paths = [
            "C:/Windows/Fonts/simhei.ttf",  # Windows 黑体
            "C:/Windows/Fonts/msyh.ttc",    # Windows 微软雅黑
            "C:/Windows/Fonts/simsun.ttc",  # Windows 宋体
            "/System/Library/Fonts/PingFang.ttc",  # macOS
            "/System/Library/Fonts/STHeiti Medium.ttc",  # macOS
            "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",  # Linux
        ]
        
        for font_path in font_paths:
            if os.path.exists(font_path):
                return font_path
        
        return None
    
    def show_wordcloud_image(self, image_path):
        """在主线程中显示词云图"""
        try:
            # 检查文件是否存在
            if not os.path.exists(image_path):
                messagebox.showerror("错误", f"无法找到图片文件: {image_path}")
                return
            
            # 根据文件格式处理
            file_ext = os.path.splitext(image_path)[1].lower()
            
            if file_ext == '.pdf':
                # PDF文件使用简单版本显示
                self.show_pdf_as_image_simple(image_path)
            else:
                # 图片文件（PNG、JPG等）
                self.show_image_file(image_path)
                
        except Exception as e:
            print(f"显示词云图时出错: {e}")
            messagebox.showerror("错误", f"无法显示词云图: {str(e)}")

    def show_image_file(self, image_path):
        """显示图片文件"""
        try:
            # 创建新窗口显示词云
            img_window = tk.Toplevel(self.root)
            img_window.title("词云图预览")
            img_window.geometry("800x600")
            
            # 加载并显示图片
            img = Image.open(image_path)
            img = img.resize((780, 580), Image.Resampling.LANCZOS)
            photo = ImageTk.PhotoImage(img)
            
            img_label = tk.Label(img_window, image=photo)
            img_label.image = photo
            img_label.pack(padx=10, pady=10)
            
            # 添加按钮框架
            button_frame = tk.Frame(img_window)
            button_frame.pack(pady=(0, 10))
            
            # 保存按钮
            save_btn = tk.Button(button_frame, text="另存为", 
                                command=lambda: self.save_image_as(img, image_path),
                                bg=self.btn_color, fg="white",
                                font=("微软雅黑", 10))
            save_btn.pack(side=tk.LEFT, padx=5)
            
            # 关闭按钮
            close_btn = tk.Button(button_frame, text="关闭", 
                                 command=img_window.destroy,
                                 bg="#7f8c8d", fg="white",
                                 font=("微软雅黑", 10))
            close_btn.pack(side=tk.LEFT, padx=5)
            
        except Exception as e:
            print(f"显示图片文件时出错: {e}")
            messagebox.showerror("错误", f"无法显示图片: {str(e)}")

    def show_pdf_as_image_simple(self, pdf_path):
        """使用简单方式显示PDF文件信息"""
        try:
            # 创建新窗口显示词云
            img_window = tk.Toplevel(self.root)
            img_window.title("词云图预览 (PDF)")
            img_window.geometry("800x600")
            
            # 添加提示信息
            info_label = tk.Label(img_window, 
                                 text="PDF文件预览：\n请查看生成的PDF文件",
                                 fg="blue", font=("微软雅黑", 12), pady=20)
            info_label.pack()
            
            # 显示PDF文件信息
            file_size = os.path.getsize(pdf_path) / 1024  # KB
            info_text = f"文件名: {os.path.basename(pdf_path)}\n"
            info_text += f"文件大小: {file_size:.1f} KB\n"
            info_text += f"保存路径: {pdf_path}\n\n"
            info_text += "注: 请使用PDF阅读器打开查看完整内容"
            
            path_label = tk.Label(img_window, text=info_text,
                                 font=("微软雅黑", 10), justify=tk.LEFT)
            path_label.pack(pady=20)
            
            # 添加按钮框架
            button_frame = tk.Frame(img_window)
            button_frame.pack(pady=20)
            
            # 打开文件按钮
            open_btn = tk.Button(button_frame, text="打开PDF文件", 
                               command=lambda: self.open_pdf_file(pdf_path),
                               bg=self.btn_color, fg="white",
                               font=("微软雅黑", 10))
            open_btn.pack(side=tk.LEFT, padx=5)
            
            # 保存按钮
            save_btn = tk.Button(button_frame, text="另存为", 
                               command=lambda: self.save_file_as(pdf_path),
                               bg=self.btn_color, fg="white",
                               font=("微软雅黑", 10))
            save_btn.pack(side=tk.LEFT, padx=5)
            
            # 关闭按钮
            close_btn = tk.Button(button_frame, text="关闭", 
                                command=img_window.destroy,
                                bg="#7f8c8d", fg="white",
                                font=("微软雅黑", 10))
            close_btn.pack(side=tk.LEFT, padx=5)
            
        except Exception as e:
            print(f"显示PDF信息时出错: {e}")
            messagebox.showerror("错误", f"无法显示PDF信息: {str(e)}")

    def open_pdf_file(self, pdf_path):
        """使用系统默认程序打开PDF文件"""
        try:
            import subprocess
            import sys
            import os
            
            if sys.platform == 'win32':
                os.startfile(pdf_path)
            elif sys.platform == 'darwin':  # macOS
                subprocess.call(['open', pdf_path])
            else:  # Linux
                subprocess.call(['xdg-open', pdf_path])
        except Exception as e:
            print(f"打开PDF文件时出错: {e}")
            messagebox.showerror("错误", f"无法打开PDF文件: {str(e)}")

    def save_image_as(self, img, source_path):
        """将词云图另存为"""
        filetypes = [
            ("PNG图片", "*.png"),
            ("JPEG图片", "*.jpg;*.jpeg"),
            ("PDF文件", "*.pdf") if HAS_PDF_SUPPORT else ("PNG图片", "*.png"),
            ("所有文件", "*.*")
        ]
        
        # 如果源文件是PDF，默认保存为图片格式
        default_ext = ".png"
        if source_path.lower().endswith('.pdf'):
            default_ext = ".png"
            initialfile = f"{os.path.splitext(os.path.basename(source_path))[0]}.png"
        else:
            initialfile = os.path.basename(source_path)
        
        save_path = filedialog.asksaveasfilename(
            title="保存词云图",
            defaultextension=default_ext,
            filetypes=filetypes,
            initialfile=initialfile
        )
        
        if save_path:
            try:
                # 获取选择的文件格式
                ext = os.path.splitext(save_path)[1].lower()
                
                if ext == '.pdf':
                    # 如果要保存为PDF，但源图不是PDF，需要转换
                    if source_path.lower().endswith('.pdf'):
                        # 源文件已经是PDF，直接复制
                        import shutil
                        shutil.copy2(source_path, save_path)
                    else:
                        # 需要从图片转换为PDF
                        messagebox.showinfo("PDF转换", 
                            "从图片转换为PDF功能需要额外处理。\n"
                            "建议先生成PDF格式的词云，然后另存为。")
                        return
                elif ext in ['.jpg', '.jpeg']:
                    # 保存为JPEG（需要转换为RGB模式）
                    if img.mode in ('RGBA', 'LA', 'P'):
                        rgb_img = img.convert('RGB')
                        rgb_img.save(save_path, quality=95)
                    else:
                        img.save(save_path, quality=95)
                else:
                    # 保存为PNG或其他格式
                    img.save(save_path)
                    
                messagebox.showinfo("保存成功", f"词云图已保存至:\n{save_path}")
            except Exception as e:
                messagebox.showerror("保存失败", f"无法保存图片: {str(e)}")

    def save_file_as(self, source_path):
        """将文件另存为（支持PDF、图片等）"""
        filetypes = [
            ("PDF文件", "*.pdf"),
            ("PNG图片", "*.png"),
            ("JPEG图片", "*.jpg;*.jpeg"),
            ("所有文件", "*.*")
        ]
        
        # 获取源文件名
        base_name = os.path.basename(source_path)
        
        save_path = filedialog.asksaveasfilename(
            title="保存文件",
            defaultextension=os.path.splitext(source_path)[1],
            filetypes=filetypes,
            initialfile=base_name
        )
        
        if save_path:
            try:
                import shutil
                shutil.copy2(source_path, save_path)
                messagebox.showinfo("保存成功", f"文件已保存至:\n{save_path}")
            except Exception as e:
                messagebox.showerror("保存失败", f"无法保存文件: {str(e)}")

def main():
    """主函数"""
    root = tk.Tk()
    app = WordFrequencyAnalyzer(root)
    root.mainloop()

if __name__ == "__main__":
    main()